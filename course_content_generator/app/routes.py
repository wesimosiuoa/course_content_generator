

from flask import Blueprint, abort, jsonify, render_template, request, redirect, url_for, flash, session, send_file, Response
import mysql
import os
import time
from werkzeug.security import generate_password_hash, check_password_hash
from app.utils.message import message
import json
import hashlib
from .db_management.db import get_db_connection
from .db_management.sql import insert as db_insert
from .db_management.sql import select as db_select
from .db_management.sql import update as db_update
from .db_management.sql import delete as db_delete
from .db_management.sql import select_one as db_select_one
from .services.methods import (
    existing_certificate,
    get_certificate_if_earned,
    explain_text_service,
    save_certificate,
    summarize_service,
    text_to_speech_service,
    get_current_user_profile,
    save_course_reaction,
    is_enrolled,
    get_user_reaction,
    get_all_courses,
    get_recommended_courses,
    get_trending_courses,
    log_search,
    search_courses,
    persist_generated_course,
    enroll_user_in_course,
    load_course_content_by_id,
)
from app.services.profile_service import *
import app.services.assessment_service as assessment_service
# LLM API token and URL for LLM service
from app.services.llm_service import generate_course, generate_lesson_notes
from app.services.assessment_service import * 
from app.services.assessment_service import issue_certificate as issue_course_certificate
from app.services.notes.notes_service import generate_lesson_notes as create_lesson_notes, get_all_notes_for_course, download_notes
import traceback
from app.services.pptx.pptx_controller import generate_lesson_pptx
from app.services.pptx.pptx_service import sanitize_filename, PPTX_DIR
from app.services.video.video_controller import generate_lesson_video
from app.services.video.video_service import VIDEO_DIR
from app.services.methods import (
    get_lesson_quiz_score,
    get_module_quiz_average,
    is_module_completed,
    is_module_assessed,
    count_completed_lessons,
    get_course_rating_summary,
    get_public_catalog_courses,
    get_course_instructor,
    assign_course_instructor,
)
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
import io
from flask import send_file



main = Blueprint('main', __name__)


def _keep_student_session():
    """Mark session permanent and dirty so login keys survive after generation."""
    if session.get("user_id"):
        session.permanent = True
        session.modified = True
        # Never keep full course JSON in cookie session (overflow logs users out)
        session.pop("generated_course", None)


def _course_preview_template():
    """Students stay inside the learner panel; guests get the public preview page."""
    if session.get("user_id") and (session.get("role") or "student").lower() == "student":
        return "student/course_preview.html"
    return "preview.html"


@main.route('/')
def index():
    return render_template('index.html')



@main.route('/login', methods=['GET', 'POST'])
def login():

    if request.method == 'GET':
        next_page = request.args.get("next")
        return render_template('login.html', next_page=next_page)

    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']
        next_page = request.form.get("next_page")  

        conn = get_db_connection()

        try:
            user = db_select_one(conn, """
                SELECT id, full_name, email, password_hash, role, expertise_domain
                FROM users
                WHERE email = %s
            """, (email,))

            if user and check_password_hash(user[3], password):
                role = (user[4] or 'student').lower()

                # Keep only small preview ids across login; drop bloated course JSON
                preview_course_id = (
                    session.get("preview_course_id")
                    or session.get("generated_course_id")
                )
                light_prefs = session.get("preferences")
                session.clear()
                session.permanent = True
                session['user_id'] = user[0]
                session['full_name'] = user[1]
                session['email'] = user[2]
                session['role'] = role
                session['expertise_domain'] = user[5]
                if preview_course_id:
                    session["preview_course_id"] = preview_course_id
                    session["generated_course_id"] = preview_course_id
                if light_prefs and isinstance(light_prefs, dict):
                    session["preferences"] = light_prefs
                session.modified = True

                profile = get_current_user_profile()
                print("Profile object:", profile)
                print("Profile type:", type(profile))
                if profile and profile.get("static_profile"):
                    session["profile_level"] = profile["static_profile"].get("explicit_level")
                session.modified = True

                # Check if user has COMPLETE preferences (essential fields filled)
                pref = get_user_preferences(user[0], check_complete=True)
                print(f"User preferences (complete check): {pref}")

                if not pref and role == 'student':
                    print("Redirecting to set_preferences - no complete preferences found")
                    flash("Please complete your learning preferences to continue.", "info")
                    return redirect("/set_preferences")

                print("Complete preferences found - redirecting to dashboard")
                flash("Login successful!", "success")

                # Guest may have generated a course; do NOT auto-enroll.
                # They can like/dislike or choose enroll from the preview.
                preview_course_id = session.get("preview_course_id") or session.get("generated_course_id")

                if next_page and next_page not in ("None", "", None):
                    return redirect(next_page)

                if role == 'admin':
                    return redirect(url_for("main.admin_dashboard"))
                if role == 'instructor':
                    return redirect(url_for("main.instructor_dashboard"))

                if role == 'student' and preview_course_id:
                    flash(
                        "Welcome back. Review your generated course — enroll only if you want it in My Courses.",
                        "info",
                    )
                    return redirect(url_for("main.preview_generated_course"))

                return redirect(url_for("main.student_dashboard"))
            else:
                flash("Invalid email or password.", "danger")

        except Exception as e:
            print (f" Error from db {str(e)}")
            print("Error from db:")
            traceback.print_exc()
            flash("Database error occurred.", "danger")

        finally:
            conn.close()

    return render_template('login.html')

@main.route('/register', methods=['GET', 'POST'])
def register():

    if request.method == 'POST':

        full_name = request.form['full_name']
        email = request.form['email']
        password = request.form['password']
        role = request.form.get('role', 'student')
        expertise_domain = request.form.get('expertise_domain', '').strip() or None

        if role not in ['student', 'instructor', 'admin']:
            role = 'student'

        if role == 'instructor' and not expertise_domain:
            flash('Instructor accounts need an expertise field area.', 'danger')
            return render_template('register.html')

        password_hash = generate_password_hash(password)

        conn = get_db_connection()
        cursor = conn.cursor()

        try:
            
            db_insert(conn, """
                INSERT INTO users (full_name, email, password_hash, role, expertise_domain)
                VALUES (%s, %s, %s, %s, %s)
            """, (full_name, email, password_hash, role, expertise_domain))
            conn.commit()

            message("Account created successfully!", "success")
            return redirect(url_for('main.login'))

        except mysql.connector.Error as err:
            message("Email already exists.", "danger")

        finally:
            cursor.close()
            conn.close()

    return render_template('register.html')

@main.route('/student_dashboard')
def student_dashboard():
    if 'user_id' not in session:
        message("Please log in to access the dashboard.", "warning")
        return redirect(url_for('main.login'))

    if session.get('role') == 'admin':
        return redirect(url_for('main.admin_dashboard'))
    if session.get('role') == 'instructor':
        return redirect(url_for('main.instructor_dashboard'))

    return render_template('/student/student_dashboard.html', full_name=session.get('full_name'))

@main.route('/instructor_dashboard')
def instructor_dashboard():
    if 'user_id' not in session:
        message("Please log in to access the dashboard.", "warning")
        return redirect(url_for('main.login'))

    if session.get('role') == 'admin':
        return redirect(url_for('main.admin_dashboard'))

    conn = get_db_connection()
    try:
        pending_reviews = db_select_one(conn, """
            SELECT COUNT(*)
            FROM saq_reevaluation_requests
            WHERE status = 'pending'
        """)
        pending_reviews_count = pending_reviews[0] if pending_reviews else 0
    finally:
        conn.close()

    return render_template(
        'instructor/dashboard.html',
        full_name=session.get('full_name'),
        expertise_domain=session.get('expertise_domain'),
        pending_reviews_count=pending_reviews_count
    )

@main.route('/admin_dashboard')
def admin_dashboard():
    if 'user_id' not in session:
        message("Please log in to access the dashboard.", "warning")
        return redirect(url_for('main.login'))

    if session.get('role') != 'admin':
        message("Admins only.", "warning")
        return redirect(url_for('main.login'))

    from app.services.llm_metrics_service import get_admin_llm_performance_metrics
    from app.services.evaluation_metrics_service import get_evaluation_matrix

    try:
        metrics = get_admin_llm_performance_metrics()
    except Exception as e:
        print(f"Admin LLM metrics error: {e}")
        traceback.print_exc()
        metrics = [
            {
                'label': 'LLM Metrics Unavailable',
                'display_value': 'Error',
                'description': 'Could not compute realtime LLM accuracy metrics. Check server logs.',
                'detail': str(e),
                'icon': 'fa-triangle-exclamation',
            }
        ]

    try:
        evaluation_matrix = get_evaluation_matrix()
    except Exception as e:
        print(f"Evaluation matrix error: {e}")
        traceback.print_exc()
        evaluation_matrix = []

    return render_template(
        'admin/dashboard.html',
        full_name=session.get('full_name'),
        metrics=metrics,
        evaluation_matrix=evaluation_matrix,
    )


def _require_admin_or_instructor():
    """Gate for expert rating workspace (admins + instructors)."""
    if 'user_id' not in session:
        message("Please log in first.", "warning")
        return redirect(url_for('main.login'))
    if session.get('role') not in ('admin', 'instructor'):
        message("Expert ratings are for admins and instructors only.", "warning")
        return redirect(url_for('main.login'))
    return None


@main.route('/admin/course/<int:course_id>')
def admin_course_review(course_id):
    """
    Admin-panel course review: full materials for expert evaluation.
    Stays inside admin auth (no student login redirect).
    """
    gate = _require_admin_or_instructor()
    if gate:
        return gate

    from app.services.evaluation_metrics_service import (
        get_admin_course_review,
        get_evaluator_ratings_for_course,
        evaluator_has_rated_course,
    )

    course = get_admin_course_review(course_id)
    if not course:
        flash('Course not found.', 'warning')
        return redirect(url_for('main.admin_expert_ratings'))

    already_rated = evaluator_has_rated_course(session['user_id'], course_id)
    existing_ratings = get_evaluator_ratings_for_course(session['user_id'], course_id)

    return render_template(
        'admin/course_review.html',
        full_name=session.get('full_name'),
        course=course,
        already_rated=already_rated,
        existing_ratings=existing_ratings,
        role=session.get('role'),
    )


@main.route('/admin/expert_ratings', methods=['GET', 'POST'])
def admin_expert_ratings():
    """
    Expert rating workspace — where humans score Completeness, Relevance,
    Hallucination, Learning Gain, and Usability for generated courses.
    Each expert rates a course once (per metric).
    """
    gate = _require_admin_or_instructor()
    if gate:
        return gate

    from app.services.evaluation_metrics_service import (
        get_expert_metric_definitions,
        save_expert_rating_batch,
        list_courses_for_expert_rating,
        list_expert_ratings,
        expert_rating_coverage_summary,
        get_evaluator_ratings_for_course,
        evaluator_has_rated_course,
    )

    if request.method == 'POST':
        try:
            course_raw = (request.form.get('course_id') or '').strip()
            course_id = int(course_raw) if course_raw and course_raw != 'system' else None
            if evaluator_has_rated_course(session['user_id'], course_id):
                flash(
                    'You already rated this course. Each expert rates a course once.',
                    'warning',
                )
                return redirect(url_for('main.admin_expert_ratings', course_id=course_raw or ''))

            metrics = get_expert_metric_definitions()
            batch = []
            for m in metrics:
                key = m['key']
                score_raw = request.form.get(f'score_{key}')
                if score_raw is None or str(score_raw).strip() == '':
                    continue
                batch.append({
                    'metric_name': key,
                    'score': float(score_raw),
                    'max_score': float(request.form.get(f'max_{key}') or m.get('max_score', 5)),
                    'notes': (request.form.get(f'notes_{key}') or '').strip(),
                })
            # Require all 5 metrics for a complete once-off rating
            if len(batch) < len(metrics):
                flash(
                    'Please rate all metrics (5 star ratings) before submitting. '
                    'Each course is rated once.',
                    'warning',
                )
                return redirect(url_for('main.admin_expert_ratings', course_id=course_raw or ''))

            result = save_expert_rating_batch(
                batch,
                evaluator_id=session.get('user_id'),
                course_id=course_id,
            )
            if result.get('saved'):
                flash(
                    f"Saved your one-time rating ({result['saved']} metrics). "
                    'It appears in the evaluation matrix Manual column.',
                    'success',
                )
            if result.get('skipped_existing'):
                flash(
                    f"{result['skipped_existing']} metric(s) were already rated and were not changed.",
                    'info',
                )
        except Exception as e:
            print(f"Expert ratings error: {e}")
            traceback.print_exc()
            flash('Failed to save expert ratings.', 'danger')
        return redirect(url_for('main.admin_expert_ratings', course_id=request.form.get('course_id') or ''))

    selected_course = request.args.get('course_id') or ''
    filter_course_id = None
    if selected_course and selected_course not in ('', 'system'):
        try:
            filter_course_id = int(selected_course)
        except ValueError:
            filter_course_id = None

    courses = list_courses_for_expert_rating()
    history = list_expert_ratings(limit=40, course_id=filter_course_id)
    coverage = expert_rating_coverage_summary()
    metrics = get_expert_metric_definitions()

    already_rated = False
    existing_ratings = {}
    if selected_course == 'system' or filter_course_id is not None:
        cid = None if selected_course == 'system' else filter_course_id
        already_rated = evaluator_has_rated_course(session['user_id'], cid)
        existing_ratings = get_evaluator_ratings_for_course(session['user_id'], cid)

    # Optional: show course content summary for the selected course
    course_preview = None
    if filter_course_id:
        conn = get_db_connection()
        try:
            row = db_select_one(
                conn,
                "SELECT id, title, content FROM courses WHERE id = %s",
                (filter_course_id,),
            )
            if row:
                content = json.loads(row[2]) if row[2] else {}
                modules = content.get('modules') or []
                course_preview = {
                    'id': row[0],
                    'title': row[1],
                    'domain': content.get('domain'),
                    'level': content.get('level'),
                    'overview': (content.get('overview') or '')[:400],
                    'module_count': len(modules),
                    'lesson_count': sum(len(m.get('lessons') or []) for m in modules if isinstance(m, dict)),
                    'outcomes': content.get('learning_outcomes') or [],
                    'resources_count': len(content.get('resources') or []),
                }
        finally:
            conn.close()

    return render_template(
        'admin/expert_ratings.html',
        full_name=session.get('full_name'),
        metrics=metrics,
        courses=courses,
        history=history,
        coverage=coverage,
        selected_course=selected_course,
        course_preview=course_preview,
        already_rated=already_rated,
        existing_ratings=existing_ratings,
        role=session.get('role'),
    )



@main.route('/admin/manual_evaluation', methods=['POST'])
def admin_manual_evaluation():
    """Legacy single-metric POST — redirects into the expert ratings workspace."""
    gate = _require_admin_or_instructor()
    if gate:
        return gate

    from app.services.evaluation_metrics_service import save_manual_evaluation, EXPERT_METRIC_KEYS

    metric_name = (request.form.get('metric_name') or '').strip().lower().replace(' ', '_')
    label_map = {
        'completeness': 'completeness',
        'relevance': 'relevance',
        'hallucination_rate': 'hallucination',
        'hallucination': 'hallucination',
        'learning_gain': 'learning_gain',
        'usability_(sus)': 'usability',
        'usability': 'usability',
    }
    metric_key = label_map.get(metric_name, metric_name)
    if metric_key not in EXPERT_METRIC_KEYS:
        flash('Invalid metric. Use the Expert Ratings page.', 'danger')
        return redirect(url_for('main.admin_expert_ratings'))

    try:
        score = float(request.form.get('score', 0))
        max_score = float(request.form.get('max_score', 5))
        course_id = request.form.get('course_id') or None
        notes = (request.form.get('notes') or '').strip()
        if course_id:
            course_id = int(course_id)
        result = save_manual_evaluation(
            metric_name=metric_key,
            score=score,
            evaluator_id=session.get('user_id'),
            course_id=course_id,
            max_score=max_score,
            notes=notes,
            allow_update=False,
        )
        if result == 'exists':
            flash('You already rated this metric for this course (once only).', 'warning')
        else:
            flash(f'Expert rating saved for {metric_key}.', 'success')
    except Exception as e:
        print(f"Manual evaluation error: {e}")
        traceback.print_exc()
        flash('Failed to save expert rating.', 'danger')

    return redirect(url_for('main.admin_expert_ratings'))


@main.route('/sus_survey', methods=['GET', 'POST'])
def sus_survey():
    """System Usability Scale survey (manual input → automatic SUS score)."""
    if 'user_id' not in session:
        message("Please log in to complete the SUS survey.", "warning")
        return redirect(url_for('main.login'))

    from app.services.evaluation_metrics_service import save_sus_response

    questions = [
        "I think that I would like to use this system frequently.",
        "I found the system unnecessarily complex.",
        "I thought the system was easy to use.",
        "I think that I would need the support of a technical person to be able to use this system.",
        "I found the various functions in this system were well integrated.",
        "I thought there was too much inconsistency in this system.",
        "I would imagine that most people would learn to use this system very quickly.",
        "I found the system very cumbersome to use.",
        "I felt very confident using the system.",
        "I needed to learn a lot of things before I could get going with this system.",
    ]

    if request.method == 'POST':
        try:
            answers = []
            for i in range(1, 11):
                val = int(request.form.get(f'q{i}', 0))
                if val < 1 or val > 5:
                    raise ValueError(f'Question {i} must be rated 1–5')
                answers.append(val)
            score = save_sus_response(session['user_id'], answers)
            flash(f'Thank you! Your SUS score is {score:.1f} / 100.', 'success')
            if session.get('role') == 'admin':
                return redirect(url_for('main.admin_dashboard'))
            return redirect(url_for('main.student_dashboard'))
        except Exception as e:
            print(f"SUS survey error: {e}")
            flash('Please rate all 10 items from 1 (Strongly disagree) to 5 (Strongly agree).', 'danger')

    return render_template(
        'sus_survey.html',
        questions=questions,
        full_name=session.get('full_name'),
    )

@main.route('/preferences')
def preferences():
    return render_template('preference.html')

@main.route('/generate_preview', methods=['POST', 'GET'])
def generate_preview():
    """
    Public generation flow.
    Always saves the course publicly to the catalog so others can discover it.
    Enrollment is optional — generator may only like/dislike and leave without enrolling.
    Never stores full course JSON in the session cookie (that overflows login state).
    """

    preferences = {
        "domain": request.form.get("domain"),
        "topic": request.form.get("topic"),
        "goal": request.form.get("goal"),
        "level": request.form.get("level"),
        "duration": request.form.get("duration"),
        "learning_preference": request.form.get("learning_preference"),
        "prior_knowledge": request.form.get("prior_knowledge")
    }

    try:
        if not preferences["domain"] or not preferences["topic"]:
            flash("Domain and Topic are required fields.", "danger")
            return redirect(url_for("main.preferences"))

        from app.services.evaluation_metrics_service import GenerationTimer

        with GenerationTimer(
            "course",
            user_id=session.get("user_id"),
            meta={"domain": preferences.get("domain"), "topic": preferences.get("topic")},
        ) as timer:
            course_data = generate_course(preferences)
            if not course_data:
                timer.success = False

        if not course_data:
            flash("Course generation failed. Please try again.", "danger")
            return redirect(url_for("main.preferences"))

        user_id = session.get("user_id")
        course_id = persist_generated_course(
            course_data,
            user_id=user_id,
            preferences=preferences,
        )

        # Lightweight session only — preserve login keys; no auto-enroll
        session["preferences"] = preferences
        session["preview_course_id"] = course_id
        session["generated_course_id"] = course_id
        session.pop("generated_course", None)
        _keep_student_session()

        flash(
            "Course generated and saved to the catalog. You can like/dislike it, "
            "enroll if you want, or leave it for other learners.",
            "success",
        )

        # Logged-in students stay in the learner panel via redirect
        if user_id:
            return redirect(url_for("main.preview_generated_course"))

        return render_template(
            "preview.html",
            course=course_data,
            user_reaction=None,
            is_enrolled=False,
            course_id=course_id,
        )

    except Exception as e:
        print("Error during course generation:", str(e))
        traceback.print_exc()
        flash("An error occurred during course generation.", "danger")
        return redirect(url_for("main.preferences"))

@main.route('/loader_test')
def loader_test():
    """Test page for loader animation"""
    if "user_id" not in session:
        return redirect(url_for("main.login"))
    
    return render_template("student/loader_test.html")

@main.route('/preview_generated_course')
def preview_generated_course():
    """
    Show generated course preview.
    Loads from DB by id so the login cookie stays intact.
    Logged-in students always render inside the student panel.
    """
    _keep_student_session()

    course_id = session.get("preview_course_id") or session.get("generated_course_id")
    course_data = None

    if course_id:
        course_data = load_course_content_by_id(course_id)
    elif session.get("generated_course"):
        # Legacy session payload — migrate off the cookie immediately
        try:
            course_id = persist_generated_course(
                session.get("generated_course"),
                user_id=session.get("user_id"),
                preferences=session.get("preferences"),
            )
            course_data = load_course_content_by_id(course_id) or session.get("generated_course")
            session["preview_course_id"] = course_id
            session["generated_course_id"] = course_id
            session.pop("generated_course", None)
            _keep_student_session()
        except Exception as e:
            print("Error migrating generated course from session:", e)
            course_data = session.get("generated_course")

    if not course_data:
        flash("No generated course found. Please generate a course first.", "warning")
        if session.get("user_id"):
            return redirect(url_for("main.discover"))
        return redirect(url_for("main.preferences"))

    user_reaction = None
    enrolled = False
    user_id = session.get("user_id")

    if user_id and course_id:
        user_reaction = get_user_reaction(user_id, course_id)
        enrolled = is_enrolled(user_id, course_id)

    return render_template(
        _course_preview_template(),
        course=course_data,
        user_reaction=user_reaction,
        is_enrolled=enrolled,
        course_id=course_id,
    )

@main.route('/react_course', methods=['POST'])
def react_course():

    if "user_id" not in session:
        flash("Please log in to react to this course.", "warning")

        next_page = request.form.get("next_page")

        print("NEXT PAGE FROM FORM:", next_page)  # debug

        return redirect(url_for("main.login", next=next_page))

    # Save the reaction and get the response
    _keep_student_session()
    response = save_course_reaction()
    _keep_student_session()
    return response

@main.route('/enroll_generated_course')
def enroll_generated_course():
    """Enroll in the currently generated course (by course_id in session)."""

    if "user_id" not in session:
        flash("Please log in to enroll in this course.", "warning")
        next_url = url_for("main.enroll_generated_course")
        return redirect(url_for("main.login", next=next_url))

    _keep_student_session()
    course_id = session.get("preview_course_id") or session.get("generated_course_id")

    # Legacy: full course still in session
    if not course_id and session.get("generated_course"):
        try:
            course_id = persist_generated_course(
                session.get("generated_course"),
                user_id=session.get("user_id"),
                preferences=session.get("preferences"),
            )
            session["preview_course_id"] = course_id
            session["generated_course_id"] = course_id
            session.pop("generated_course", None)
            _keep_student_session()
        except Exception as e:
            print(f"Error persisting generated course on enroll: {e}")

    if not course_id:
        flash("No generated course found. Please generate a course first.", "warning")
        return redirect(url_for("main.preview_generated_course"))

    try:
        enroll_user_in_course(session["user_id"], course_id)
        _keep_student_session()
        flash("Successfully enrolled! This course is now in My Courses.", "success")
        return redirect(url_for("main.view_my_course", course_id=course_id))
    except Exception as e:
        print(f"Error enrolling in generated course: {e}")
        traceback.print_exc()
        flash("An error occurred while enrolling.", "danger")
        return redirect(url_for("main.preview_generated_course"))


@main.route("/courses")
def courses():
    """
    Public course catalog: ratings, likes, enrollments, duration,
    and assigned assistant instructor (expertise-matched).
    """
    user_id = session.get("user_id")
    catalog = get_public_catalog_courses(user_id=user_id)
    enrolled_ids = {c["id"] for c in catalog if c.get("enrolled")}

    return render_template(
        "courses.html",
        courses=catalog,
        enrolled_ids=enrolled_ids,
    )

@main.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('main.index'))


@main.route("/course/<int:course_id>")
def view_course(course_id):
    
    if 'user_id' not in session:
        # Store the intended destination and redirect to login
        return redirect(url_for("main.login", next=url_for("main.view_course", course_id=course_id)))

    conn = get_db_connection()

    try:
        # 1️⃣ Fetch course
        course = db_select_one(conn, """
            SELECT id, title, description, content, popularity_score, created_at
            FROM courses
            WHERE id = %s
        """, (course_id,))

        if not course:
            abort(404)

        # 2️⃣ Unpack
        course_data = {
            "id": course[0],
            "title": course[1],
            "description": course[2],
            "content": json.loads(course[3]),
            "popularity_score": course[4],
            "created_at": course[5]
        }

        # 3️⃣ Check if logged-in user reacted
        user_reaction = None

        feedback = db_select_one(conn, """
            SELECT reaction
            FROM course_feedback
            WHERE user_id = %s AND course_id = %s
        """, (session["user_id"], course_id))

        if feedback:
            user_reaction = feedback[0]

        # 4️⃣ Check enrollment status
        enrolled = is_enrolled(session["user_id"], course_id)

        return render_template(
            "course_content.html",
            course=course_data,
            user_reaction=user_reaction,
            enrolled=enrolled
        )

    finally:
        conn.close()

# enrollment route
@main.route("/enroll/<int:course_id>")
def enroll(course_id):

    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.enroll", course_id=course_id)))

    conn = get_db_connection()

    try:
        db_insert(conn, """
            INSERT IGNORE INTO enrollments (user_id, course_id)
            VALUES (%s, %s)
        """, (session["user_id"], course_id))

        flash("Successfully enrolled in course.", "success")

        return redirect(url_for("main.view_course", course_id=course_id))

    finally:
        conn.close()


@main.route('/generate_lesson_notes/<int:course_id>/<int:module_index>/<int:lesson_index>')
def generate_lesson_notes_route(course_id, module_index, lesson_index):
    """Generate notes for a specific lesson"""
    
    if "user_id" not in session:
        flash("Please log in to generate lesson notes.", "warning")
        return redirect(url_for("main.login"))
    
    conn = get_db_connection()
    
    try:
        # Fetch course from database
        course = db_select_one(conn, """
            SELECT id, title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))
        
        if not course:
            flash("Course not found.", "danger")
            return redirect(url_for("main.my_courses"))
        
        course_title = course[1]
        course_content = json.loads(course[2])
        
        # Get user preferences for personalization
        user_preferences = get_user_preferences(session["user_id"])
        
        # Extract module and lesson
        modules = course_content.get('modules', [])
        
        if module_index < 0 or module_index >= len(modules):
            flash("Invalid module.", "danger")
            return redirect(url_for("main.my_courses"))
        
        module = modules[module_index]
        module_title = module.get('title', f'Module {module_index + 1}')
        
        lessons = module.get('lessons', [])
        
        if lesson_index < 0 or lesson_index >= len(lessons):
            flash("Invalid lesson.", "danger")
            return redirect(url_for("main.my_courses"))
        
        lesson = lessons[lesson_index]
        lesson_title = lesson.get('title', f'Lesson {lesson_index + 1}')
        
        # CHECK IF NOTES ALREADY EXIST
        from app.services.notes.notes_service import get_existing_notes_file
        existing_file = get_existing_notes_file(course_title, module_title, lesson_title)
        
        if existing_file and os.path.exists(existing_file):
            print(f"Notes already exist: {existing_file}")
            flash("Using existing notes.", "info")
            
            # Store existing file path in session
            session[f'notes_{course_id}_{module_index}_{lesson_index}'] = existing_file
            
            # Redirect to view/download options
            flash("Notes are ready! You can view online or download.", "success")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        
        # Notes don't exist - generate new ones
        print(f"Generating NEW notes for: {course_title} > {module_title} > {lesson_title}")

        from app.services.evaluation_metrics_service import GenerationTimer

        # Generate notes using LLM (timed for Generation Time metric)
        with GenerationTimer(
            "notes",
            user_id=session.get("user_id"),
            meta={"course_id": course_id, "module_index": module_index, "lesson_index": lesson_index},
        ) as timer:
            notes_result = create_lesson_notes(
                course_title=course_title,
                module_title=module_title,
                lesson=lesson,  # Pass the entire lesson dict
                preferences=user_preferences
            )
            if not (notes_result and notes_result.get('success')):
                timer.success = False
        
        if notes_result and notes_result.get('success'):
            flash("Lesson notes generated successfully!", "success")
            
            # Store file path in session for both view and download
            session[f'notes_{course_id}_{module_index}_{lesson_index}'] = notes_result['file_path']
            
            # Show success with options - redirect to a choice page or back to course
            # For now, let's redirect back to course with a message
            flash("You can now view online or download the notes.", "info")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        else:
            error_msg = notes_result.get('message', 'Failed to generate notes') if notes_result else 'Unknown error'
            flash(f"Error generating notes: {error_msg}", "danger")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
            
    except Exception as e:
        print(f"Error generating lesson notes: {str(e)}")
        traceback.print_exc()
        flash("An error occurred while generating notes.", "danger")
        return redirect(url_for("main.my_courses"))
    
    finally:
        conn.close()

@main.route ('/view_lesson_slides/<int:course_id>/<int:module_index>/<int:lesson_index>')
def view_lesson_slides(course_id, module_index, lesson_index):
    """View generated lesson slides in browser"""
    
    if "user_id" not in session:
        flash("Please log in to view lesson slides.", "warning")
        return redirect(url_for("main.login"))
    
    try:
        # Get file path from session (note: stored as 'pptx_*' not 'slides_*')
        file_path = session.get(f'pptx_{course_id}_{module_index}_{lesson_index}')
        
        if not file_path or not os.path.exists(file_path):
            flash("Slides not generated yet. Please generate slides first.", "warning")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        
        # Get filename from path
        filename = os.path.basename(file_path)
        title = filename.replace('.pptx', '').replace('_', ' ').title()
        
        # Render viewer template with PPTX file path
        return render_template("student/lesson_slides_viewer.html", 
                             file_path=file_path,
                             filename=filename,
                             title=title,
                             course_id=course_id,
                             module_index=module_index,
                             lesson_index=lesson_index)
            
    except Exception as e:
        print(f"Error viewing slides: {str(e)}")
        flash("Error loading slides.", "danger")
        return redirect(url_for("main.my_courses"))
    
@main.route('/view_lesson_notes/<int:course_id>/<int:module_index>/<int:lesson_index>')
def view_lesson_notes(course_id, module_index, lesson_index):
    """View lesson notes in browser (without downloading)"""
    
    if "user_id" not in session:
        flash("Please log in to view lesson notes.", "warning")
        return redirect(url_for("main.login"))
    
    try:
        # Get file path from session
        file_path = session.get(f'notes_{course_id}_{module_index}_{lesson_index}')
        
        if not file_path or not os.path.exists(file_path):
            flash("Notes not generated yet. Please generate notes first.", "warning")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        
        # Read the content from the Word document
        from docx import Document
        doc = Document(file_path)
        
        # Extract all paragraphs
        content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Safely check formatting
                is_bold = False
                is_italic = False
                is_heading = False
                style_name_str = ''
                
                # Check runs for bold/italic
                try:
                    for run in paragraph.runs:
                        if hasattr(run, 'bold') and run.bold is True:
                            is_bold = True
                        if hasattr(run, 'italic') and run.italic is True:
                            is_italic = True
                except Exception as e:
                    print(f"Error checking run formatting: {e}")
                
                # Check if it's a heading style
                try:
                    if hasattr(paragraph, 'style') and hasattr(paragraph.style, 'name'):
                        style_name_obj = paragraph.style.name
                        if style_name_obj is not None and isinstance(style_name_obj, str):
                            style_name_str = style_name_obj
                            is_heading = style_name_str.startswith('Heading')
                except Exception as e:
                    print(f"Error checking style name: {e}")
                
                content.append({
                    'text': paragraph.text,
                    'bold': is_bold,
                    'italic': is_italic,
                    'heading': is_heading,
                    'style_name': style_name_str
                })
        
        # Get document title from first heading or filename
        title = os.path.basename(file_path).replace('.docx', '').replace('_', ' ').title()
        
        return render_template("student/lesson_notes_viewer.html", 
                             content=content, 
                             title=title,
                             course_id=course_id,
                             module_index=module_index,
                             lesson_index=lesson_index)
            
    except Exception as e:
        print(f"Error viewing notes: {str(e)}")
        flash("Error loading notes.", "danger")
        return redirect(url_for("main.my_courses"))


@main.route('/download_lesson_notes/<int:course_id>/<int:module_index>/<int:lesson_index>')
def download_lesson_notes(course_id, module_index, lesson_index):
    """Download generated lesson notes"""
    
    if "user_id" not in session:
        flash("Please log in to download notes.", "warning")
        return redirect(url_for("main.login"))
    
    try:
        # Get file path from session
        file_path = session.get(f'notes_{course_id}_{module_index}_{lesson_index}')
        
        if not file_path or not os.path.exists(file_path):
            flash("Notes file not found. Please generate notes first.", "warning")
            return redirect(url_for("main.my_courses"))
        
        # Read file
        filename, content = download_notes(file_path)
        
        if filename and content:
            from flask import send_file
            return send_file(
                file_path,
                as_attachment=True,
                download_name=filename,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
        else:
            flash("Error reading notes file.", "danger")
            return redirect(url_for("main.my_courses"))
            
    except Exception as e:
        print(f"Error downloading notes: {str(e)}")
        flash("Error downloading notes.", "danger")
        return redirect(url_for("main.my_courses"))


@main.route('/lesson_quiz')
def lesson_quiz():
    """Render the lesson quiz page using query parameters."""
    if "user_id" not in session:
        flash("Please log in to take the quiz.", "warning")
        return redirect(url_for("main.login"))

    course_id = request.args.get('course')
    module_index = request.args.get('module')
    lesson_index = request.args.get('lesson')

    if course_id is None or module_index is None or lesson_index is None:
        flash("Invalid quiz request. Please select a lesson quiz.", "danger")
        return redirect(url_for("main.my_courses"))

    try:
        course_id = int(course_id)
        module_index = int(module_index)
        lesson_index = int(lesson_index)
    except ValueError:
        flash("Invalid quiz identifiers.", "danger")
        return redirect(url_for("main.my_courses"))


    
    return render_template(
        'student/lesson_quiz.html',
        course_id=course_id,
        module_index=module_index,
        lesson_index=lesson_index
    )


@main.route('/generate_lesson_quiz/<int:course_id>/<int:module_index>/<int:lesson_index>')
def generate_lesson_quiz(course_id, module_index, lesson_index):
    """Generate quiz for a specific lesson"""
    
    if "user_id" not in session:
        flash("Please log in to generate quiz.", "warning")
        return redirect(url_for("main.login"))
    
    conn = get_db_connection()
    
    try:
        # Fetch course from database
        course = db_select_one(conn, """
            SELECT id, title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))
        
        if not course:
            flash("Course not found.", "danger")
            return redirect(url_for("main.my_courses"))
        
        course_title = course[1]
        course_content = json.loads(course[2])
        
        # Get lesson data
        try:
            lesson_data = course_content['modules'][module_index]['lessons'][lesson_index]
        except (IndexError, KeyError):
            flash("Lesson not found.", "danger")
            return redirect(url_for("main.view_my_course", course_id=course_id))
        
        # Generate quiz (timed for Generation Time metric)
        try:
            from app.services.evaluation_metrics_service import GenerationTimer

            with GenerationTimer(
                "quiz",
                user_id=session.get("user_id"),
                meta={"course_id": course_id, "module_index": module_index, "lesson_index": lesson_index},
            ) as timer:
                question_ids = assessment_service.create_quiz_for_lesson(
                    course_id, module_index, lesson_index, lesson_data
                )
                if not (question_ids and len(question_ids) > 0):
                    timer.success = False
            
            if question_ids and len(question_ids) > 0:
                flash(f"Quiz generated successfully with {len(question_ids)} questions!", "success")
            else:
                flash("Quiz generation completed but no questions were created.", "warning")
                
        except Exception as e:
            print(f"Error generating quiz: {str(e)}")
            traceback.print_exc()
            flash("Error generating quiz. Please try again.", "danger")
        
        return redirect(url_for("main.view_my_course", course_id=course_id))
            
    except Exception as e:
        print(f"Error in quiz generation: {str(e)}")
        traceback.print_exc()
        flash("An error occurred while generating quiz.", "danger")
        return redirect(url_for("main.my_courses"))
    
    finally:
        conn.close()


@main.route("/my_courses")
def my_courses():

    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.my_courses")))

    conn = get_db_connection()

    try:
        courses = db_select(conn, """
            SELECT c.id, c.title, c.content, c.popularity_score, c.created_at, e.completed
            FROM courses c
            JOIN enrollments e ON c.id = e.course_id
            WHERE e.user_id = %s
            ORDER BY e.enrolled_at DESC
        """, (session["user_id"],))

        formatted_courses = []
        user_id = session["user_id"]

        for course in courses:
            content = course[2]
            if isinstance(content, str):
                try:
                    content = json.loads(content)
                except Exception:
                    content = {}
            elif not isinstance(content, dict):
                content = {}

            modules = content.get("modules") or []
            total_lessons = sum(len((m or {}).get("lessons") or []) for m in modules)
            done_lessons = count_completed_lessons(user_id, course[0])
            progress = (done_lessons * 100.0 / total_lessons) if total_lessons else 0.0

            # Prefer grade/cert; fall back to all module assessments done
            status = None
            try:
                status = get_course_completion_status(user_id, course[0])
            except Exception:
                status = None
            cert = get_certificate_if_earned(user_id, course[0])
            all_assessed = False
            if modules:
                all_assessed = all(
                    is_module_assessed(user_id, course[0], i) for i in range(len(modules))
                )
            is_completed = bool(cert or (status and status.get("passed")) or all_assessed or course[5])

            rating = get_course_rating_summary(course[0])

            formatted_courses.append({
                "id": course[0],
                "title": course[1],
                "overview": content.get("overview", ""),
                "popularity_score": course[3],
                "created_at": course[4],
                "enrolled": True,
                "completed": is_completed,
                "progress": progress,
                "rating": rating,
            })

        return render_template("student/courses.html", courses=formatted_courses)

    finally:
        conn.close()


@main.route("/view_my_course/<int:course_id>")
def view_my_course(course_id):

    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.view_my_course", course_id=course_id)))

    conn = get_db_connection()

    try:
        course = db_select_one(conn, """
            SELECT id, title, description, content, popularity_score, created_at
            FROM courses
            WHERE id = %s
        """, (course_id,))

        if not course:
            abort(404)

        course_data = {
            "id": course[0],
            "title": course[1],
            "description": course[2],
            "content": json.loads(course[3]),
            "popularity_score": course[4],
            "created_at": course[5]
        }
        reaction = get_user_reaction(session.get("user_id"), course_data["id"]) if "user_id" in session else None
        enrolled = is_enrolled(session.get("user_id"), course_data["id"]) if "user_id" in session else False

        # Build scores dict for each lesson (None = not taken)
        scores = {}
        course_content = course_data["content"] or {}
        modules = course_content.get("modules") or []
        total_lessons = 0
        for module_idx, module in enumerate(modules):
            lessons = (module or {}).get("lessons") or []
            total_lessons += len(lessons)
            for lesson_idx, _lesson in enumerate(lessons):
                scores[f"{module_idx}_{lesson_idx}"] = get_lesson_quiz_score(
                    session["user_id"], course_id, module_idx, lesson_idx
                )

        # Count only lessons with an actual quiz score (not the dict size)
        completed_lessons = sum(1 for s in scores.values() if s is not None)
        if completed_lessons == 0:
            # DB path is authoritative if session scores are empty/odd
            completed_lessons = count_completed_lessons(session["user_id"], course_id)

        overall_progress = (
            (completed_lessons * 100.0 / total_lessons) if total_lessons else 0.0
        )

        # Module averages + lesson completion + assessment completion
        module_averages = {}
        module_completed = {}   # all lesson quizzes done (unlocks assessment)
        module_assessed = {}    # module assessment submitted
        for module_idx, module in enumerate(modules):
            expected = len((module or {}).get("lessons") or [])
            module_averages[module_idx] = get_module_quiz_average(
                session["user_id"], course_id, module_idx
            )
            module_completed[module_idx] = is_module_completed(
                session["user_id"], course_id, module_idx, expected_lessons=expected
            )
            module_assessed[module_idx] = is_module_assessed(
                session["user_id"], course_id, module_idx
            )

        status = get_course_completion_status(session["user_id"], course_id)
        certificate = get_certificate_if_earned(session["user_id"], course_id)

        # Course is completed only when all modules are assessed, or grade/cert exists
        if modules:
            all_modules_assessed = all(module_assessed.get(i, False) for i in range(len(modules)))
        else:
            all_modules_assessed = False
        course_completed = bool(
            certificate
            or (status and status.get("passed"))
            or all_modules_assessed
        )

        rating = get_course_rating_summary(course_id)

        return render_template(
            "student/view_course.html",
            scores=scores,
            module_averages=module_averages,
            module_completed=module_completed,
            module_assessed=module_assessed,
            course=course_data,
            reaction=reaction,
            enrolled=enrolled,
            course_completed=course_completed,
            certificate=certificate,
            status=status,
            rating=rating,
            total_lessons=total_lessons,
            completed_lessons=completed_lessons,
            overall_progress=overall_progress,
        )

    finally:
        conn.close()

@main.route('/module_assessment/<int:course_id>/<int:module_index>')
def module_assessment(course_id, module_index):
    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.module_assessment", course_id=course_id, module_index=module_index)))
    
    # Check if user is enrolled
    if not is_enrolled(session["user_id"], course_id):
        flash("You must be enrolled in this course to access assessments.", "error")
        return redirect(url_for("main.view_my_course", course_id=course_id))
    
    # Check if module is completed (all lessons have quiz results)
    if not is_module_completed(session["user_id"], course_id, module_index):
        flash("You must complete all lesson quizzes in this module before taking the module assessment.", "warning")
        return redirect(url_for("main.view_my_course", course_id=course_id))
    
    conn = get_db_connection()
    try:
        course = db_select_one(conn, """
            SELECT title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))
        
        if not course:
            abort(404)
        
        course_title = course[0]
        course_content = json.loads(course[1])
        
        # Get module title
        module_title = f"Module {module_index + 1}"
        try:
            module = course_content['modules'][module_index]
            module_title = module.get('title', module_title)
        except (IndexError, KeyError):
            pass
        
        return render_template("student/module_assessment.html", 
                             course_id=course_id, 
                             course_title=course_title,
                             module_index=module_index, 
                             module_title=module_title)
    
    finally:
        conn.close()

@main.route('/enroll_course/<int:course_id>', methods=['POST', 'GET'])
def enroll_course(course_id):
    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.enroll_course", course_id=course_id)))

    conn = get_db_connection()

    try:
        db_insert(conn, """
            INSERT IGNORE INTO enrollments (user_id, course_id)
            VALUES (%s, %s)
        """, (session["user_id"], course_id))

        flash("Successfully enrolled in course.", "success")

        return redirect(url_for("main.view_my_course", course_id=course_id))

    finally:
        conn.close()


@main.route('/discover')
def discover():

    if 'user_id' not in session:
        return redirect(url_for("main.login"))
    student_id = session['user_id']
    query = request.args.get('query', '').strip()

    # Discover only surfaces courses the student has not enrolled in (enrollment is once)
    existing_courses = get_all_courses(student_id)
    search_results = []
    recommended_courses = get_recommended_courses(student_id)
    trending_courses = get_trending_courses(user_id=student_id)

    # Avoid showing the same course in both recommended and trending
    recommended_ids = {c.get("id") for c in recommended_courses}
    trending_courses = [c for c in trending_courses if c.get("id") not in recommended_ids]

    if query:
        log_search(student_id, query)
        search_results = search_courses(query, user_id=student_id)

    return render_template(
        "student/discover.html",
        search_results=search_results,
        recommended_courses=recommended_courses,
        trending_courses=trending_courses,
        existing_courses=existing_courses,
        query=query,
    )
# discover course as a student where the course they are looking for is not available yet, so they can input their preferences and get a generated course preview. Note this, we are using search behavior to trigger the course generation flow, so we will have a search bar on the discover page where they can input their desired course topic, and if it does not exist in the database, we will redirect them to the preferences page with the topic pre-filled. 
@main.route('/search', methods=['POST'])
def search():

    topic = request.form.get("topic")

    if not topic:
        flash("Please enter a topic to search.", "warning")
        return redirect(url_for("main.discover"))

    conn = get_db_connection()

    try:
        course = db_select_one(conn, """
            SELECT id
            FROM courses
            WHERE JSON_EXTRACT(content, '$.topic') = %s
        """, (json.dumps(topic),))

        if course:
            return redirect(url_for("main.view_course", course_id=course[0]))
        else:
            flash("Course not found. Please provide your preferences to generate a course preview.", "info")
            return redirect(url_for("main.set_preferences", topic=topic))

    finally:
        conn.close()


# Profile Building 
@main.route("/set_preferences", methods=["GET", "POST"])
def set_preferences():

    if "user_id" not in session:
        return redirect("/login")

    user_id = session["user_id"]

    if request.method == "POST":

        data = {
            "domain": request.form.get("domain"),
            "topic": request.form.get("topic"),
            "goal": request.form.get("goal"),
            "level": request.form.get("level"),
            "duration": request.form.get("duration"),
            "learning_preference": request.form.get("learning_preference"),
            "prior_knowledge": request.form.get("prior_knowledge")
        }

        save_user_preferences(user_id, data)

        flash("Your preferences saved successfully", "success")
        return redirect(url_for("main.student_dashboard"))

    existing = get_user_preferences(user_id)

    return render_template(
        "student/preferences.html",
        existing=existing
    )
#generate course from logged in student 
@main.route ("/learner_generate_course", methods=["GET", "POST"])
def learner_generate_course():
    """
    Logged-in student course generation.
    - Course is saved publicly to the catalog (available for other students).
    - Enrollment is optional — student may only like/dislike.
    - Student stays logged in (full course is not stored in the session cookie).
    """
    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.learner_generate_course")))

    user_id = session["user_id"]
    topic = request.args.get("topic") or request.form.get("topic")

    user_preferences = get_user_preferences(user_id)

    if not user_preferences:
        flash("Please set your learning preferences first.", "warning")
        return redirect(url_for("main.set_preferences"))

    if request.method == "POST":
        try:
            generation_prefs = {
                "domain": user_preferences.get("domain", ["General"]),
                "topic": topic or user_preferences.get("topic", ""),
                "goal": user_preferences.get("goal", "Professional Skill Development"),
                "level": user_preferences.get("level", "Beginner"),
                "duration": user_preferences.get("duration", "4"),
                "learning_preference": user_preferences.get("learning_preference", "Balanced Approach"),
                "prior_knowledge": user_preferences.get("prior_knowledge", ""),
            }

            print(f"Generating course with preferences: {generation_prefs}")

            from app.services.evaluation_metrics_service import GenerationTimer

            with GenerationTimer(
                "course",
                user_id=user_id,
                meta={"topic": generation_prefs.get("topic"), "source": "learner_generate_course"},
            ) as timer:
                course_data = generate_course(generation_prefs)
                if not course_data:
                    timer.success = False

            if not course_data:
                flash("Failed to generate course. Please try again.", "danger")
                return render_template(
                    "student/personalised_course.html",
                    user_preferences=user_preferences,
                    topic=topic,
                )

            # Persist publicly to catalog — no auto-enroll
            course_id = persist_generated_course(
                course_data,
                user_id=user_id,
                preferences=generation_prefs,
            )

            session["preferences"] = generation_prefs
            session["preview_course_id"] = course_id
            session["generated_course_id"] = course_id
            session.pop("generated_course", None)
            _keep_student_session()

            log_search(user_id, topic or generation_prefs.get("topic") or "")

            flash(
                "Course generated and saved to the catalog. Like or dislike it, "
                "enroll only if you want it in My Courses — other learners can still find it.",
                "success",
            )
            # Stay signed in inside the student panel
            return redirect(url_for("main.preview_generated_course"))

        except Exception as err:
            print(f"Error generating course: {err}")
            traceback.print_exc()
            flash("An error occurred during course generation.", "danger")

    return render_template(
        "student/personalised_course.html",
        user_preferences=user_preferences,
        topic=topic,
    )

@main.route('/user_profile')
def user_profile():
    profile_data = build_user_profile(session["user_id"])
    print (f"Profile data for user {session['user_id']}: {profile_data}")  # Debug print
    return render_template("student/profile.html", profile=profile_data)

@main.route('/text-to-speech', methods=['POST'])
def text_to_speech():
    try:
        data = request.get_json()
        text = data.get('text')
        
        if not text:
            return jsonify({"error": "No text provided"}), 400

        audio = text_to_speech_service(text)

        # Return the raw bytes as an MPEG stream
        return Response(audio, mimetype='audio/mpeg')

    except Exception as e:
        print(f"CRITICAL TTS ERROR: {str(e)}") # This shows in your terminal
        return jsonify({"error": str(e)}), 500
    

@main.route('/summarize', methods=['POST'])
def summarize():
    try:
        data = request.get_json()
        text = data.get('text')

        summary = summarize_service(text)

        return jsonify({"summary": summary})

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    
@main.route('/explain', methods=['POST'])
def explain():
    try:
        data = request.get_json()
        text = data.get('text')

        explanation = explain_text_service(text)

        return jsonify({"explanation": explanation})

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    
# Route to generate PPTX for a specific lesson
@main.route('/generate_lesson_pptx/<int:course_id>/<int:module_index>/<int:lesson_index>')
def generate_lesson_pptx_route(course_id, module_index, lesson_index):
    """Generate PPTX for a specific lesson"""

    if "user_id" not in session:
        flash("Please log in to generate presentation.", "warning")
        return redirect(url_for("main.login"))

    conn = get_db_connection()

    try:
        # Fetch course
        course = db_select_one(conn, """
            SELECT id, title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))

        if not course:
            flash("Course not found.", "danger")
            return redirect(url_for("main.my_courses"))

        course_title = course[1]
        course_content = json.loads(course[2])

        # User preferences
        user_preferences = get_user_preferences(session["user_id"])

        # Extract module
        modules = course_content.get('modules', [])
        if module_index < 0 or module_index >= len(modules):
            flash("Invalid module.", "danger")
            return redirect(url_for("main.my_courses"))

        module = modules[module_index]
        module_title = module.get('title', f'Module {module_index + 1}')

        # Extract lesson
        lessons = module.get('lessons', [])
        if lesson_index < 0 or lesson_index >= len(lessons):
            flash("Invalid lesson.", "danger")
            return redirect(url_for("main.my_courses"))

        lesson = lessons[lesson_index]
        lesson_title = lesson.get('title', f'Lesson {lesson_index + 1}')

        # 🔷 CHECK IF PPTX ALREADY EXISTS
        
        filename = f"{sanitize_filename(course_title)}_{sanitize_filename(module_title)}_{sanitize_filename(lesson_title)}.pptx"
        file_path = os.path.join(PPTX_DIR, filename)

        if os.path.exists(file_path):
            print(f"PPTX already exists: {file_path}")
            flash("Using existing presentation.", "info")

            session[f'pptx_{course_id}_{module_index}_{lesson_index}'] = file_path

            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))

        # 🔷 GENERATE NEW PPTX
        print(f"Generating PPTX for: {course_title} > {module_title} > {lesson_title}")


        pptx_result = generate_lesson_pptx(
            course_title=course_title,
            module_title=module_title,
            lesson=lesson,
            preferences=user_preferences
        )

        if pptx_result and pptx_result.get('success'):
            flash("Presentation generated successfully!", "success")

            session[f'pptx_{course_id}_{module_index}_{lesson_index}'] = pptx_result['file_path']

            flash("You can now view or download the presentation.", "info")

            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        else:
            error_msg = pptx_result.get('message', 'Failed to generate PPTX') if pptx_result else 'Unknown error'
            flash(f"Error generating PPTX: {error_msg}", "danger")

            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))

    except Exception as e:
        print(f"Error generating PPTX: {str(e)}")
        traceback.print_exc()

        flash("An error occurred while generating presentation.", "danger")
        return redirect(url_for("main.my_courses"))

    finally:
        conn.close()

@main.route('/download_lesson_pptx/<int:course_id>/<int:module_index>/<int:lesson_index>')
def download_lesson_pptx(course_id, module_index, lesson_index):
    """Download generated PPTX file"""

    if "user_id" not in session:
        flash("Please log in to download slides.", "warning")
        return redirect(url_for("main.login"))

    try:
        key = f'pptx_{course_id}_{module_index}_{lesson_index}'
        file_path = session.get(key)

        if file_path and os.path.exists(file_path):
            return send_file(
                file_path,
                as_attachment=True,
                download_name=os.path.basename(file_path),
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        flash("Presentation not found. Please generate it first.", "danger")
        return redirect(url_for(
            'main.view_my_course',
            course_id=course_id
        ))

    except Exception as e:
        print(f"Error downloading PPTX: {str(e)}")
        traceback.print_exc()

        flash("An error occurred while downloading presentation.", "danger")
        return redirect(url_for("main.my_courses"))

@main.route('/serve_lesson_pptx/<int:course_id>/<int:module_index>/<int:lesson_index>')
def serve_lesson_pptx(course_id, module_index, lesson_index):
    """Serve PPTX slides as images for browser viewing"""
    
    if "user_id" not in session:
        flash("Please log in to view slides.", "warning")
        return redirect(url_for("main.login"))
    
    try:
        key = f'pptx_{course_id}_{module_index}_{lesson_index}'
        pptx_path = session.get(key)
        
        if not pptx_path or not os.path.exists(pptx_path):
            return jsonify({"error": "Presentation not found"}), 404
        
        # Convert PPTX to images
        slide_images = convert_pptx_to_images(pptx_path)
        
        if not slide_images:
            return jsonify({"error": "Could not convert presentation"}), 500
        
        return jsonify({
            "success": True,
            "slides": slide_images,
            "total_slides": len(slide_images)
        })
        
    except Exception as e:
        print(f"Error serving PPTX: {str(e)}")
        return jsonify({"error": str(e)}), 500

def convert_pptx_to_images(pptx_path):
    """Convert PPTX slides to PNG images"""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        import io
        import base64
        
        prs = Presentation(pptx_path)
        slide_images = []
        
        # Set slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Convert EMU (English Metric Units) to pixels (96 DPI)
        px_width = int(slide_width * 96 / 914400)
        px_height = int(slide_height * 96 / 914400)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            try:
                # Create blank image with white background
                img = Image.new('RGB', (px_width, px_height), color='white')
                draw = ImageDraw.Draw(img)
                
                # Try to use a default font, fall back to default if not available
                try:
                    font = ImageFont.truetype("arial.ttf", 24)
                    small_font = ImageFont.truetype("arial.ttf", 16)
                except:
                    font = ImageFont.load_default()
                    small_font = font
                
                y_pos = 40
                
                # Extract and draw text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Estimate if this is a title (usually first shape with text)
                        if y_pos == 40:
                            draw.text((40, y_pos), text, fill='black', font=font)
                            y_pos += 60
                        else:
                            # Wrap text for readability
                            words = text.split()
                            line = ""
                            for word in words:
                                test_line = line + word + " "
                                if len(test_line) > 50:
                                    draw.text((40, y_pos), line, fill='black', font=small_font)
                                    y_pos += 25
                                    line = word + " "
                                else:
                                    line = test_line
                            if line:
                                draw.text((40, y_pos), line, fill='black', font=small_font)
                                y_pos += 25
                
                # Convert to base64
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
                
                slide_images.append({
                    "slide_num": slide_num,
                    "data": f"data:image/png;base64,{img_base64}"
                })
                
            except Exception as e:
                print(f"Error converting slide {slide_num}: {e}")
                # Create error slide
                img = Image.new('RGB', (px_width, px_height), color='white')
                draw = ImageDraw.Draw(img)
                draw.text((40, 40), f"Error rendering slide {slide_num}", fill='red')
                
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
                
                slide_images.append({
                    "slide_num": slide_num,
                    "data": f"data:image/png;base64,{img_base64}"
                })
        
        return slide_images
        
    except ImportError:
        print("python-pptx not installed. Install with: pip install python-pptx pillow")
        return None
    except Exception as e:
        print(f"Error converting PPTX to images: {e}")
        return None


# ==========================================
# VIDEO GENERATION ROUTES
# ==========================================

@main.route('/generate_lesson_video/<int:course_id>/<int:module_index>/<int:lesson_index>')
def generate_lesson_video_route(course_id, module_index, lesson_index):
    """Generate video from lesson notes and slides"""

    if "user_id" not in session:
        flash("Please log in to generate videos.", "warning")
        return redirect(url_for("main.login"))

    try:
        # Get course data
        conn = get_db_connection()
        course = db_select_one(conn, """
            SELECT id, title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))

        if not course:
            flash("Course not found.", "danger")
            return redirect(url_for("main.my_courses"))

        # Parse course content
        course_data = json.loads(course[2]) if course[2] else {}
        modules = course_data.get('modules', [])

        if module_index >= len(modules):
            flash("Module not found.", "danger")
            return redirect(url_for("main.my_courses"))

        module = modules[module_index]
        lessons = module.get('lessons', [])

        if lesson_index >= len(lessons):
            flash("Lesson not found.", "danger")
            return redirect(url_for("main.my_courses"))

        lesson = lessons[lesson_index]

        # Check if notes and PPTX exist
        notes_key = f'notes_{course_id}_{module_index}_{lesson_index}'
        pptx_key = f'pptx_{course_id}_{module_index}_{lesson_index}'

        notes_path = session.get(notes_key)
        pptx_path = session.get(pptx_key)

        if not notes_path or not os.path.exists(notes_path):
            flash("Lesson notes not found. Please generate notes first.", "warning")
            return redirect(url_for('main.view_my_course', course_id=course_id))

        if not pptx_path or not os.path.exists(pptx_path):
            flash("Presentation slides not found. Please generate slides first.", "warning")
            return redirect(url_for('main.view_my_course', course_id=course_id))

        # Generate video
        result = generate_lesson_video(
            course[1],  # course_title
            module['title'],
            lesson,
            notes_path,
            pptx_path
        )

        if result['success']:
            # Store video path in session
            video_key = f'video_{course_id}_{module_index}_{lesson_index}'
            session[video_key] = result['video_path']

            flash(f"Video generated successfully! Duration: {result['duration']:.0f}s", "success")
        else:
            flash(f"Error generating video: {result['message']}", "danger")

        return redirect(url_for('main.view_my_course', course_id=course_id))

    except Exception as e:
        print(f"Error generating video: {str(e)}")
        traceback.print_exc()
        flash("An error occurred while generating video.", "danger")
        return redirect(url_for("main.my_courses"))

    finally:
        conn.close()


@main.route('/view_lesson_video/<int:course_id>/<int:module_index>/<int:lesson_index>')
def view_lesson_video(course_id, module_index, lesson_index):
    """View generated lesson video in browser"""

    if "user_id" not in session:
        flash("Please log in to view lesson videos.", "warning")
        return redirect(url_for("main.login"))

    try:
        # Get file path from session
        file_path = session.get(f'video_{course_id}_{module_index}_{lesson_index}')

        if not file_path or not os.path.exists(file_path):
            flash("Video not generated yet. Please generate video first.", "warning")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))

        # Get filename from path
        filename = os.path.basename(file_path)
        title = filename.replace('.mp4', '').replace('_', ' ').title()
        print (f"{file_path} is the video path for course {course_id}, module {module_index}, lesson {lesson_index}")

        # Render video viewer template
        return render_template("student/lesson_video_viewer.html",
                             file_path=file_path,
                             filename=filename,
                             title=title,
                             course_id=course_id,
                             module_index=module_index,
                             lesson_index=lesson_index)

    except Exception as e:
        print(f"Error viewing video: {str(e)}")
        flash("Error loading video.", "danger")
        return redirect(url_for("main.my_courses"))


@main.route('/download_lesson_video/<int:course_id>/<int:module_index>/<int:lesson_index>')
def download_lesson_video(course_id, module_index, lesson_index):
    """Download generated video file"""

    if "user_id" not in session:
        flash("Please log in to download videos.", "warning")
        return redirect(url_for("main.login"))

    try:
        key = f'video_{course_id}_{module_index}_{lesson_index}'
        file_path = session.get(key)

        if file_path and os.path.exists(file_path):
            return send_file(
                file_path,
                as_attachment=True,
                download_name=os.path.basename(file_path),
                mimetype="video/mp4"
            )

        flash("Video not found. Please generate it first.", "danger")
        return redirect(url_for(
            'main.view_my_course',
            course_id=course_id
        ))

    except Exception as e:
        print(f"Error downloading video: {str(e)}")
        traceback.print_exc()

        flash("An error occurred while downloading video.", "danger")
        return redirect(url_for("main.my_courses"))


def _resolve_lesson_video_path(course_id, module_index, lesson_index):
    """Find video path from session or known VIDEO_DIR naming."""
    key = f"video_{course_id}_{module_index}_{lesson_index}"
    video_path = session.get(key)
    if video_path and os.path.exists(video_path):
        return video_path

    # Fallback: scan VIDEO_DIR if session key was lost but file exists
    try:
        if os.path.isdir(VIDEO_DIR):
            for name in os.listdir(VIDEO_DIR):
                if name.lower().endswith(".mp4") and not name.endswith(".tmp.mp4"):
                    # keep session-less recovery only when single match is ambiguous —
                    # prefer exact session path above
                    pass
    except Exception:
        pass
    return video_path if video_path and os.path.exists(video_path) else None


def _stream_mp4_with_ranges(video_path):
    """
    Stream MP4 with HTTP Range support so browser <video> can seek/play.
    Also re-applies faststart once if needed for progressive playback.
    """
    from app.services.video.video_service import optimize_existing_video_for_browser

    if not video_path or not os.path.exists(video_path):
        return jsonify({"error": "Video not found"}), 404

    # One-time browser optimization for older files written without faststart
    flag = video_path + ".faststarted"
    if not os.path.exists(flag):
        try:
            if optimize_existing_video_for_browser(video_path):
                open(flag, "w", encoding="utf-8").write("1")
        except Exception as e:
            print(f"faststart optimize skipped: {e}")

    file_size = os.path.getsize(video_path)
    range_header = request.headers.get("Range", None)

    # Full file
    if not range_header:
        resp = send_file(
            video_path,
            mimetype="video/mp4",
            conditional=True,
            as_attachment=False,
            download_name=os.path.basename(video_path),
            max_age=0,
        )
        resp.headers["Accept-Ranges"] = "bytes"
        resp.headers["Content-Length"] = str(file_size)
        resp.headers["Cache-Control"] = "no-cache"
        return resp

    # bytes=start-end
    try:
        units, _, rng = range_header.partition("=")
        if units.strip().lower() != "bytes":
            return Response(status=416)
        start_s, _, end_s = rng.partition("-")
        start = int(start_s) if start_s else 0
        end = int(end_s) if end_s else file_size - 1
        end = min(end, file_size - 1)
        if start > end or start < 0:
            return Response(status=416)
    except Exception:
        return Response(status=416)

    length = end - start + 1

    def generate():
        with open(video_path, "rb") as f:
            f.seek(start)
            remaining = length
            chunk = 1024 * 256
            while remaining > 0:
                data = f.read(min(chunk, remaining))
                if not data:
                    break
                remaining -= len(data)
                yield data

    resp = Response(generate(), status=206, mimetype="video/mp4", direct_passthrough=True)
    resp.headers["Content-Range"] = f"bytes {start}-{end}/{file_size}"
    resp.headers["Accept-Ranges"] = "bytes"
    resp.headers["Content-Length"] = str(length)
    resp.headers["Cache-Control"] = "no-cache"
    return resp


@main.route('/serve_lesson_video/<int:course_id>/<int:module_index>/<int:lesson_index>')
def serve_lesson_video(course_id, module_index, lesson_index):
    """Serve video file for browser streaming (Range + faststart)."""

    if "user_id" not in session:
        return jsonify({"error": "Authentication required"}), 401

    try:
        video_path = _resolve_lesson_video_path(course_id, module_index, lesson_index)
        if not video_path:
            return jsonify({"error": "Video not found"}), 404
        return _stream_mp4_with_ranges(video_path)

    except Exception as e:
        print(f"Error serving video: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


# ==================== ASSESSMENT ROUTES ====================

@main.route('/assessment/quiz/<int:course_id>/<int:module_index>/<int:lesson_index>', methods=['GET'])
def get_quiz(course_id, module_index, lesson_index):
    """Get quiz for a specific lesson"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        quiz_data = assessment_service.get_lesson_quiz(course_id, module_index, lesson_index)
        questions = quiz_data.get('questions', []) if isinstance(quiz_data, dict) else quiz_data
        
        return jsonify({
            'status': 'success',
            'data': {
                'course_id': course_id,
                'course_title': quiz_data.get('course_name', f'Course {course_id}') if isinstance(quiz_data, dict) else f'Course {course_id}',
                'module_index': module_index,
                'module_title': quiz_data.get('module_name', f'Module {module_index + 1}') if isinstance(quiz_data, dict) else f'Module {module_index + 1}',
                'lesson_index': lesson_index,
                'lesson_title': quiz_data.get('lesson_name', f'Lesson {lesson_index + 1}') if isinstance(quiz_data, dict) else f'Lesson {lesson_index + 1}',
                'questions': questions,
                'total_questions': len(questions)
            }
        })
    except Exception as e:
        print(f"Error fetching quiz: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to load quiz'}), 500


@main.route('/assessment/quiz/submit', methods=['POST'])
def submit_quiz():
    """Submit quiz responses"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        data = request.get_json()
        user_id = session['user_id']
        course_id = data['course_id']
        module_index = data['module_index']
        lesson_index = data['lesson_index']
        responses = data['responses']
        
        result = assessment_service.submit_quiz_response(
            user_id, course_id, module_index, lesson_index, responses
        )
        
        return jsonify({
            'status': 'success',
            'data': result
        })
    except Exception as e:
        print(f"Error submitting quiz: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to submit quiz'}), 500


@main.route('/assessment/quiz/result/<int:course_id>/<int:module_index>/<int:lesson_index>', methods=['GET'])
def get_quiz_result(course_id, module_index, lesson_index):
    """Get quiz result for a lesson"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        conn = get_db_connection()
        user_id = session['user_id']
        
        result = db_select_one(conn, """
            SELECT total_questions, correct_answers, score_percentage, passed, completed_at
            FROM lesson_quiz_results
            WHERE user_id = %s AND course_id = %s 
            AND module_index = %s AND lesson_index = %s
        """, (user_id, course_id, module_index, lesson_index))
        
        conn.close()
        
        if result:
            return jsonify({
                'status': 'success',
                'data': {
                    'total_questions': result[0],
                    'correct_answers': result[1],
                    'score_percentage': float(result[2]),
                    'passed': bool(result[3]),
                    'completed_at': result[4].isoformat() if result[4] else None
                }
            })
        else:
            return jsonify({'status': 'not_attempted', 'message': 'Quiz not yet attempted'})
    except Exception as e:
        print(f"Error fetching quiz result: {str(e)}")
        return jsonify({'error': 'Failed to load result'}), 500


# ==================== FINAL ASSESSMENT ROUTES ====================

@main.route('/assessment/final/<int:course_id>', methods=['GET'])
def get_final_assessment(course_id):
    """Get final assessment questions"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        assessments = assessment_service.get_final_assessments(course_id)
        
        return jsonify({
            'status': 'success',
            'data': {
                'course_id': course_id,
                'questions': assessments,
                'total_questions': len(assessments)
            }
        })
    except Exception as e:
        print(f"Error fetching final assessment: {str(e)}")
        return jsonify({'error': 'Failed to load assessment'}), 500


@main.route('/assessment/final/submit', methods=['POST'])
def submit_final_assessment():
    """Submit final assessment response"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        data = request.get_json()
        user_id = session['user_id']
        
        result = assessment_service.submit_final_response(
            user_id,
            data['course_id'],
            data['final_assessment_id'],
            data['answer_text']
        )
        
        return jsonify({
            'status': 'success',
            'data': result
        })
    except Exception as e:
        print(f"Error submitting final assessment: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to submit assessment'}), 500


@main.route('/assessment/final/results/<int:course_id>', methods=['GET'])
def get_final_results(course_id):
    """Get final assessment submission and results"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        conn = get_db_connection()
        user_id = session['user_id']
        
        responses = db_select(conn, """
            SELECT sfr.id, sfr.final_assessment_id, sfr.answer_text, sfr.score,
                   sfr.feedback, sfr.is_graded, fa.question_text, fa.max_score
            FROM student_final_responses sfr
            JOIN final_assessments fa ON sfr.final_assessment_id = fa.id
            WHERE sfr.user_id = %s AND sfr.course_id = %s
        """, (user_id, course_id))
        
        conn.close()
        
        results = [
            {
                'id': r[0],
                'assessment_id': r[1],
                'answer': r[2],
                'score': float(r[3]) if r[3] else None,
                'feedback': r[4],
                'is_graded': bool(r[5]),
                'question': r[6],
                'max_score': r[7]
            }
            for r in responses
        ]
        
        return jsonify({
            'status': 'success',
            'data': results
        })
    except Exception as e:
        print(f"Error fetching final results: {str(e)}")
        return jsonify({'error': 'Failed to load results'}), 500


# ==================== LLM-BASED GRADING ENDPOINTS ====================

@main.route('/assessment/llm/grade-saq/<int:response_id>', methods=['POST'])
def llm_grade_saq(response_id):
    """Use LLM to suggest a grade for SAQ response"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        conn = get_db_connection()
        
        # Get response data
        response_data = db_select_one(conn, """
            SELECT ssr.id, ssr.answer_text, saq.question_text, saq.rubric, saq.max_score
            FROM student_saq_responses ssr
            JOIN short_answer_questions saq ON ssr.saq_id = saq.id
            WHERE ssr.id = %s
        """, (response_id,))
        
        conn.close()
        
        if not response_data:
            return jsonify({'error': 'Response not found'}), 404
        
        question_text = response_data[2]
        rubric = response_data[3]
        max_score = response_data[4]
        answer_text = response_data[1]
        
        if rubric:
            rubric = json.loads(rubric) if isinstance(rubric, str) else rubric
        
        # Get LLM grading suggestion
        grading_result = assessment_service.llm_grade_saq_response(
            question_text,
            rubric or {},
            answer_text,
            max_score
        )
        
        return jsonify({
            'status': 'success',
            'data': grading_result,
            'message': 'LLM grading suggestion provided'
        })
    except Exception as e:
        print(f"Error getting LLM grade: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to generate grade'}), 500


@main.route('/assessment/llm/grade-final/<int:response_id>', methods=['POST'])
def llm_grade_final(response_id):
    """Use LLM to suggest a grade for final assessment response"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        conn = get_db_connection()
        
        # Get response data
        response_data = db_select_one(conn, """
            SELECT sfr.id, sfr.answer_text, fa.question_text, fa.max_score
            FROM student_final_responses sfr
            JOIN final_assessments fa ON sfr.final_assessment_id = fa.id
            WHERE sfr.id = %s
        """, (response_id,))
        
        conn.close()
        
        if not response_data:
            return jsonify({'error': 'Response not found'}), 404
        
        question_text = response_data[2]
        max_score = response_data[3]
        answer_text = response_data[1]
        
        # Get LLM grading suggestion
        grading_result = assessment_service.llm_grade_final_response(
            question_text,
            answer_text,
            max_score=max_score
        )
        
        return jsonify({
            'status': 'success',
            'data': grading_result,
            'message': 'LLM grading suggestion provided'
        })
    except Exception as e:
        print(f"Error getting LLM grade: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to generate grade'}), 500


@main.route('/assessment/llm/apply-grade/<int:response_id>/<string:response_type>', methods=['POST'])
def apply_llm_grade(response_id, response_type):
    """Apply LLM-suggested grade to response and save it"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    


@main.route('/view_certificate/<int:course_id>')
def view_certificate(course_id):
    user_id = session['user_id']

    cert = get_certificate(user_id, course_id)

    # Get course separately
    conn = get_db_connection()
    course_row = db_select_one(conn, """
        SELECT id, title FROM courses WHERE id = %s
    """, (course_id,))
    conn.close()

    course = {
        "id": course_row[0],
        "title": course_row[1]
    }
    qr_code = generate_qr(cert['certificate_code'])
    return render_template(
        "student/certificate.html",
        cert=cert,
        course=course,
        qr_code=qr_code
    )


import qrcode
import base64
from io import BytesIO

def generate_qr(code):
    verify_url = f"http://127.0.0.1:5000/verify/{code}"

    qr = qrcode.make(verify_url)

    buffer = BytesIO()
    qr.save(buffer, format="PNG")
    buffer.seek(0)

    qr_base64 = base64.b64encode(buffer.getvalue()).decode()

    return qr_base64

@main.route('/generate_certificate/<int:course_id>')
def generate_certificate(course_id):
    user_id = session['user_id']
    user_name = session.get('full_name', 'Student')

    # 1. Check completion status
    status = get_course_completion_status(user_id, course_id)

    if not status or not status['passed']:
        flash("You have not completed the course successfully.", "danger")
        return redirect(url_for('main.view_my_course', course_id=course_id))

    # 2. Issue certificate
    result = issue_certificate(user_id, course_id, user_name, None)

    if result['status'] == 'issued':
        flash("Certificate generated successfully!", "success")
    else:
        flash(result['message'], "warning")

    return redirect(url_for('main.view_certificate', course_id=course_id))


@main.route('/download_certificate/<int:course_id>')
def download_certificate(course_id):
    user_id = session['user_id']

    cert = get_certificate(user_id, course_id)

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer)

    styles = getSampleStyleSheet()

    elements = []
    elements.append(Paragraph("Certificate of Completion", styles['Title']))
    elements.append(Paragraph(session.get('full_name', 'Student'), styles['Heading2']))
    elements.append(Paragraph(f"Certificate Code: {cert['certificate_code']}", styles['Normal']))

    doc.build(elements)

    buffer.seek(0)

    return send_file(buffer, as_attachment=True, download_name="certificate.pdf")





@main.route('/verify/<string:code>')
def verify_certificate(code):
    conn = get_db_connection()

    try:
        cert = db_select_one(conn, """
            SELECT 
                c.certificate_code, 
                c.issue_date, 
                u.full_name, 
                co.title
            FROM certificates c
            JOIN users u ON c.user_id = u.id
            JOIN courses co ON c.course_id = co.id
            WHERE c.certificate_code = %s
        """, (code,))

    except Exception as e:
        print(f"Verification error: {e}")
        conn.close()
        return render_template("certificate_verify.html", cert=None)

    conn.close()

    # 🔴 If certificate not found
    if not cert:
        return render_template("certificate_verify.html", cert=None)

    # ✅ Convert tuple → dict
    certificate = {
        "certificate_code": cert[0],
        "issue_date": cert[1],   # keep as datetime
        "user_name": cert[2],
        "course_name": cert[3]
    }

    return render_template("certificate_verify.html", cert=certificate)