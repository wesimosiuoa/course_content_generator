"""
Video Generation Service
Generates MP4 videos from lesson notes and slides.
TTS order: ElevenLabs → gTTS → Windows SAPI → silent audio (always playable).
"""

import os
import re
import json
import wave
import math
import struct
import tempfile
import shutil
import subprocess
from datetime import datetime
from pathlib import Path
from docx import Document
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import moviepy.editor as mpy
from dotenv import load_dotenv

load_dotenv()

print("API Key loaded:", bool(os.getenv("ELEVEN_API_KEY")))

ELEVENLABS_API_KEY = os.getenv("ELEVEN_API_KEY", "")
ELEVENLABS_VOICE_ID = os.getenv("VOICE_ID", "EXAVITQu4vr4xnSDxMaL")

VIDEO_DIR = os.path.join(os.path.dirname(__file__), "videos")
MAX_VIDEO_DURATION = 300  # 5 minutes
TARGET_FPS = 24
MIN_SLIDE_SECONDS = 4.0
WORDS_PER_SECOND = 2.4  # approx speech rate for timing estimates


def ensure_video_directory():
    if not os.path.exists(VIDEO_DIR):
        os.makedirs(VIDEO_DIR)
    return VIDEO_DIR


def _ffmpeg_executable():
    """Prefer system ffmpeg; fall back to imageio-ffmpeg (bundled with moviepy)."""
    exe = shutil.which("ffmpeg")
    if exe:
        return exe
    try:
        import imageio_ffmpeg
        return imageio_ffmpeg.get_ffmpeg_exe()
    except Exception:
        return None


def _ensure_mp4_faststart(src_path, dest_path=None):
    """
    Remux MP4 so the moov atom is at the start (required for browser <video>).
    Returns True if dest_path is ready for progressive browser playback.
    """
    dest_path = dest_path or src_path
    if not src_path or not os.path.exists(src_path):
        return False

    ffmpeg = _ffmpeg_executable()
    if not ffmpeg:
        print("No ffmpeg available for faststart remux")
        if os.path.abspath(src_path) != os.path.abspath(dest_path):
            try:
                shutil.copy2(src_path, dest_path)
                return True
            except Exception:
                return False
        return os.path.exists(dest_path)

    # Remux with stream copy + faststart
    try:
        tmp = dest_path + ".faststart.mp4"
        if os.path.exists(tmp):
            try:
                os.remove(tmp)
            except Exception:
                pass
        cmd = [
            ffmpeg, "-y",
            "-i", src_path,
            "-c", "copy",
            "-movflags", "+faststart",
            tmp,
        ]
        completed = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=180,
        )
        if completed.returncode == 0 and os.path.exists(tmp) and os.path.getsize(tmp) > 0:
            if os.path.exists(dest_path) and os.path.abspath(tmp) != os.path.abspath(dest_path):
                try:
                    os.remove(dest_path)
                except Exception:
                    pass
            os.replace(tmp, dest_path)
            print(f"✓ MP4 faststart applied: {dest_path}")
            return True
        print("ffmpeg faststart failed:", (completed.stderr or "")[-500:])
        try:
            if os.path.exists(tmp):
                os.remove(tmp)
        except Exception:
            pass
    except Exception as e:
        print(f"faststart remux error: {e}")

    if os.path.abspath(src_path) != os.path.abspath(dest_path):
        try:
            shutil.copy2(src_path, dest_path)
            return os.path.exists(dest_path)
        except Exception as e:
            print(f"copy after faststart fail: {e}")
            return False
    return os.path.exists(dest_path)


def optimize_existing_video_for_browser(video_path):
    """Re-faststart an already-generated MP4 so the browser can stream it."""
    if not video_path or not os.path.exists(video_path):
        return False
    try:
        fixed = video_path + ".web.mp4"
        ok = _ensure_mp4_faststart(video_path, fixed)
        if ok and os.path.exists(fixed):
            os.replace(fixed, video_path)
            return True
    except Exception as e:
        print(f"optimize_existing_video_for_browser: {e}")
    return False


def sanitize_filename(text):
    invalid_chars = '<>:"/\\|?*'
    for char in invalid_chars:
        text = text.replace(char, "_")
    return text.replace(" ", "_").lower()[:50]


def extract_text_from_docx(docx_path):
    try:
        doc = Document(docx_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text.strip())
        # tables
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    text.append(" — ".join(cells))
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""


def extract_slides_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        slides = []

        for slide_num, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": slide_num + 1,
                "title": "",
                "text_content": "",
            }
            text_parts = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                raw = (shape.text or "").strip()
                if not raw:
                    continue
                # Prefer shorter first text as title
                if not slide_data["title"] and len(raw) < 120 and "\n" not in raw:
                    slide_data["title"] = raw
                else:
                    text_parts.append(raw)

            slide_data["text_content"] = "\n".join(text_parts).strip()
            if not slide_data["title"] and text_parts:
                first_line = text_parts[0].split("\n")[0][:100]
                slide_data["title"] = first_line
            slides.append(slide_data)

        return slides
    except Exception as e:
        print(f"Error extracting slides from PPTX: {e}")
        return []


def _load_font(size=32, bold=False):
    candidates = [
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/segoeui.ttf",
        "arial.ttf",
        "DejaVuSans.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            continue
    return ImageFont.load_default()


def convert_pptx_to_images(pptx_path, output_dir):
    """Render PPTX text content onto slide images for the video."""
    try:
        prs = Presentation(pptx_path)
        image_paths = []
        title_font = _load_font(48)
        body_font = _load_font(30)

        for slide_num, slide in enumerate(prs.slides):
            img = Image.new("RGB", (1920, 1080), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)

            # Header bar
            draw.rectangle([0, 0, 1920, 110], fill=(25, 103, 210))

            y_position = 140
            max_width = 1760
            line_height = 42
            first_shape = True

            for shape in slide.shapes:
                if not hasattr(shape, "text") or not (shape.text or "").strip():
                    continue
                text = shape.text.strip()
                font = title_font if first_shape else body_font
                fill = (255, 255, 255) if first_shape and y_position < 120 else (30, 30, 30)

                # Title goes in header
                if first_shape:
                    # wrap title in header
                    words = text.replace("\n", " ").split()
                    lines, current = [], []
                    for word in words:
                        test = " ".join(current + [word])
                        bbox = draw.textbbox((0, 0), test, font=title_font)
                        if bbox[2] - bbox[0] <= 1700:
                            current.append(word)
                        else:
                            if current:
                                lines.append(" ".join(current))
                            current = [word]
                    if current:
                        lines.append(" ".join(current))
                    ty = 30
                    for line in lines[:2]:
                        draw.text((60, ty), line, fill=(255, 255, 255), font=title_font)
                        ty += 48
                    first_shape = False
                    continue

                first_shape = False
                for paragraph in text.split("\n"):
                    paragraph = paragraph.strip()
                    if not paragraph:
                        y_position += 12
                        continue
                    words = paragraph.split()
                    current_line = []
                    for word in words:
                        test_line = " ".join(current_line + [word])
                        bbox = draw.textbbox((0, 0), test_line, font=body_font)
                        if bbox[2] - bbox[0] <= max_width:
                            current_line.append(word)
                        else:
                            if current_line:
                                if y_position > 1000:
                                    break
                                draw.text((80, y_position), " ".join(current_line), fill=(40, 40, 40), font=body_font)
                                y_position += line_height
                            current_line = [word]
                    if current_line and y_position <= 1000:
                        draw.text((80, y_position), " ".join(current_line), fill=(40, 40, 40), font=body_font)
                        y_position += line_height
                    y_position += 10
                    if y_position > 1000:
                        break

            # Footer slide number
            draw.text((1780, 1020), f"{slide_num + 1}", fill=(100, 100, 100), font=_load_font(22))

            img_path = os.path.join(output_dir, f"slide_{slide_num:03d}.png")
            img.save(img_path, "PNG")
            image_paths.append(img_path)

        return image_paths
    except Exception as e:
        print(f"Error converting PPTX to images: {e}")
        return []


def _estimate_speech_seconds(text):
    words = len(re.findall(r"\w+", text or ""))
    return max(MIN_SLIDE_SECONDS, min(45.0, words / WORDS_PER_SECOND + 1.0))


def _write_silent_wav(path, duration_sec):
    """Stdlib silent WAV so video always has an audio track."""
    duration_sec = max(MIN_SLIDE_SECONDS, float(duration_sec or MIN_SLIDE_SECONDS))
    sample_rate = 22050
    n_samples = int(sample_rate * duration_sec)
    with wave.open(path, "w") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(sample_rate)
        silence = struct.pack("<h", 0)
        wf.writeframes(silence * n_samples)
    return path, duration_sec


def _tts_elevenlabs(text, output_path):
    if not ELEVENLABS_API_KEY:
        return None, 0
    try:
        from elevenlabs.client import ElevenLabs
        from elevenlabs import VoiceSettings

        client = ElevenLabs(api_key=ELEVENLABS_API_KEY)
        audio = client.text_to_speech.convert(
            text=text[:4500],
            voice_id=ELEVENLABS_VOICE_ID,
            model_id="eleven_flash_v2_5",
            voice_settings=VoiceSettings(
                stability=0.5,
                similarity_boost=0.75,
                style=0.0,
                use_speaker_boost=True,
            ),
        )
        # may be mp3 bytes
        mp3_path = output_path if output_path.lower().endswith(".mp3") else output_path + ".mp3"
        with open(mp3_path, "wb") as f:
            for chunk in audio:
                f.write(chunk)
        clip = mpy.AudioFileClip(mp3_path)
        duration = float(clip.duration or 0)
        clip.close()
        if duration <= 0:
            return None, 0
        return mp3_path, duration
    except Exception as e:
        print(f"ElevenLabs TTS failed: {e}")
        return None, 0


def _tts_gtts(text, output_path):
    try:
        from gtts import gTTS

        mp3_path = output_path if output_path.lower().endswith(".mp3") else output_path + ".mp3"
        tts = gTTS(text=text[:4500], lang="en")
        tts.save(mp3_path)
        clip = mpy.AudioFileClip(mp3_path)
        duration = float(clip.duration or 0)
        clip.close()
        if duration <= 0:
            return None, 0
        return mp3_path, duration
    except Exception as e:
        print(f"gTTS fallback failed: {e}")
        return None, 0


def _tts_windows_sapi(text, output_path):
    """Windows built-in SAPI (no API key). Produces WAV."""
    if os.name != "nt":
        return None, 0
    try:
        wav_path = output_path if output_path.lower().endswith(".wav") else output_path + ".wav"
        # Escape for PowerShell single-quoted string
        safe = (text or "Lesson content.").replace("'", "''")[:2500]
        out_ps = wav_path.replace("'", "''")
        script = (
            "Add-Type -AssemblyName System.Speech; "
            "$s = New-Object System.Speech.Synthesis.SpeechSynthesizer; "
            f"$s.SetOutputToWaveFile('{out_ps}'); "
            f"$s.Speak('{safe}'); "
            "$s.Dispose();"
        )
        completed = subprocess.run(
            ["powershell", "-NoProfile", "-Command", script],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if completed.returncode != 0 or not os.path.exists(wav_path):
            print("Windows SAPI TTS failed:", completed.stderr[:300] if completed.stderr else "no file")
            return None, 0
        clip = mpy.AudioFileClip(wav_path)
        duration = float(clip.duration or 0)
        clip.close()
        if duration <= 0:
            return None, 0
        return wav_path, duration
    except Exception as e:
        print(f"Windows SAPI TTS failed: {e}")
        return None, 0


def generate_tts_audio(text, output_path, voice_id=None):
    """
    Generate narration audio with resilient fallbacks.
    Always returns (path, duration) when possible; silent WAV last resort.
    """
    text = (text or "").strip()
    if not text:
        text = "This section covers key ideas from the lesson."

    # 1) ElevenLabs (preferred when credits work)
    path, duration = _tts_elevenlabs(text, output_path)
    if path:
        print(f"TTS via ElevenLabs: {duration:.2f}s")
        return path, duration, "elevenlabs"

    # 2) gTTS free online fallback
    path, duration = _tts_gtts(text, output_path)
    if path:
        print(f"TTS via gTTS fallback: {duration:.2f}s")
        return path, duration, "gtts"

    # 3) Windows SAPI offline
    path, duration = _tts_windows_sapi(text, output_path)
    if path:
        print(f"TTS via Windows SAPI: {duration:.2f}s")
        return path, duration, "sapi"

    # 4) Silent track timed to estimated speech length — video still plays
    est = _estimate_speech_seconds(text)
    wav_path = output_path if output_path.lower().endswith(".wav") else output_path + ".wav"
    path, duration = _write_silent_wav(wav_path, est)
    print(f"TTS silent fallback: {duration:.2f}s (no speech service available)")
    return path, duration, "silent"


def segment_text_for_video(text, num_segments=None):
    text = (text or "").strip()
    if not text:
        return ["This lesson introduces the main ideas and applications of the topic."]

    # Prefer sentence boundaries
    sentences = re.split(r"(?<=[.!?])\s+", text)
    sentences = [s.strip() for s in sentences if s and s.strip()]
    if not sentences:
        sentences = [text]

    if not num_segments or num_segments < 1:
        num_segments = max(1, min(len(sentences), 8))

    # Distribute sentences evenly
    segments = []
    n = len(sentences)
    for i in range(num_segments):
        start = math.floor(i * n / num_segments)
        end = math.floor((i + 1) * n / num_segments)
        chunk = sentences[start:end] or [sentences[min(i, n - 1)]]
        segments.append(" ".join(chunk))

    return [s.strip() for s in segments if s.strip()]


def _slide_narration_text(slide, fallback_segment, lesson_title=""):
    title = (slide.get("title") or "").strip()
    body = (slide.get("text_content") or "").strip()
    if body and len(body.split()) >= 8:
        return f"{title}. {body}" if title and title not in body else body
    if fallback_segment:
        return f"{title}. {fallback_segment}" if title else fallback_segment
    return title or f"Key points for {lesson_title or 'this lesson'}."


def generate_video(lesson_notes_path, pptx_path, output_path, course_info=None):
    """
    Generate a playable MP4 from notes + slides.
    Continues even when ElevenLabs credits fail.
    """
    temp_dirs = []
    result = {
        "success": False,
        "video_path": None,
        "duration": 0,
        "file_size": 0,
        "slides_count": 0,
        "error": None,
        "message": "",
        "tts_provider": None,
    }

    try:
        ensure_video_directory()
        course_info = course_info or {}

        if not output_path:
            lesson_title = course_info.get("lesson_title") or "lesson"
            course_title = course_info.get("course_title") or "course"
            module_title = course_info.get("module_title") or "module"
            output_path = os.path.join(
                VIDEO_DIR,
                f"{sanitize_filename(course_title)}_{sanitize_filename(module_title)}_{sanitize_filename(lesson_title)}.mp4",
            )

        images_dir = tempfile.mkdtemp(prefix="vid_img_")
        audio_dir = tempfile.mkdtemp(prefix="vid_aud_")
        temp_dirs.extend([images_dir, audio_dir])

        print("Extracting lesson notes...")
        lesson_text = extract_text_from_docx(lesson_notes_path) if lesson_notes_path else ""

        print("Extracting slides...")
        slides = extract_slides_from_pptx(pptx_path)
        if not slides:
            raise Exception("No slides found in presentation")

        num_slides = len(slides)
        result["slides_count"] = num_slides

        if not lesson_text:
            # Build narration from slide text when notes are empty
            lesson_text = "\n".join(
                f"{s.get('title', '')}. {s.get('text_content', '')}".strip() for s in slides
            )
            if not lesson_text.strip():
                lesson_text = (
                    f"This lesson covers {course_info.get('lesson_title', 'the topic')} "
                    f"in {course_info.get('module_title', 'the module')}."
                )

        print("Converting slides to images...")
        image_paths = convert_pptx_to_images(pptx_path, images_dir)
        if not image_paths:
            image_paths = []
            for i, s in enumerate(slides):
                img = Image.new("RGB", (1920, 1080), color=(245, 248, 252))
                draw = ImageDraw.Draw(img)
                draw.rectangle([0, 0, 1920, 110], fill=(25, 103, 210))
                draw.text((60, 35), (s.get("title") or f"Slide {i+1}")[:80], fill="white", font=_load_font(40))
                body = (s.get("text_content") or lesson_text[:400])[:500]
                draw.text((80, 160), body[:900], fill=(40, 40, 40), font=_load_font(28))
                img_path = os.path.join(images_dir, f"slide_{i:03d}.png")
                img.save(img_path)
                image_paths.append(img_path)

        print("Generating audio narration (with fallbacks)...")
        text_segments = segment_text_for_video(lesson_text, num_slides)
        while len(text_segments) < num_slides:
            text_segments.append(text_segments[-1] if text_segments else lesson_text[:400])

        audio_paths = []
        segment_durations = []
        providers_used = set()

        for i in range(num_slides):
            narr = _slide_narration_text(
                slides[i] if i < len(slides) else {},
                text_segments[i] if i < len(text_segments) else lesson_text[:400],
                course_info.get("lesson_title", ""),
            )
            audio_path = os.path.join(audio_dir, f"segment_{i:03d}")
            seg_path, duration, provider = generate_tts_audio(narr, audio_path)
            providers_used.add(provider)
            if not seg_path or duration <= 0:
                wav_path = os.path.join(audio_dir, f"segment_{i:03d}_silent.wav")
                seg_path, duration = _write_silent_wav(wav_path, _estimate_speech_seconds(narr))
                providers_used.add("silent")
            audio_paths.append(seg_path)
            segment_durations.append(max(MIN_SLIDE_SECONDS, float(duration)))
            print(f"  Segment {i+1}: {duration:.2f}s via {provider}")

        result["tts_provider"] = ",".join(sorted(providers_used))
        total_audio_duration = sum(segment_durations)
        print(f"Total audio duration: {total_audio_duration:.2f}s providers={result['tts_provider']}")

        # Cap extremely long videos
        if total_audio_duration > MAX_VIDEO_DURATION and total_audio_duration > 0:
            scale = MAX_VIDEO_DURATION / total_audio_duration
            segment_durations = [max(MIN_SLIDE_SECONDS, d * scale) for d in segment_durations]

        print("Creating video clips...")
        video_clips = []
        for i, img_path in enumerate(image_paths):
            duration = segment_durations[i] if i < len(segment_durations) else MIN_SLIDE_SECONDS
            clip = mpy.ImageClip(img_path).set_duration(duration)
            if i < len(audio_paths) and audio_paths[i] and os.path.exists(audio_paths[i]):
                try:
                    audio_clip = mpy.AudioFileClip(audio_paths[i])
                    # Match audio length to slide duration
                    if audio_clip.duration and audio_clip.duration > duration:
                        audio_clip = audio_clip.subclip(0, duration)
                    clip = clip.set_audio(audio_clip)
                except Exception as audio_err:
                    print(f"  Warning: could not attach audio for slide {i+1}: {audio_err}")
            video_clips.append(clip)

        print("Assembling video...")
        final_video = mpy.concatenate_videoclips(video_clips, method="compose")
        video_duration = float(final_video.duration or 0)

        print(f"Rendering video to {output_path}...")
        os.makedirs(os.path.dirname(output_path) or VIDEO_DIR, exist_ok=True)

        # Write to a temp file first, then faststart so browsers can stream
        # (moov atom at file start). Without this, Explorer/VLC play fine but
        # <video> in Chrome/Edge often spins forever.
        tmp_out = output_path + ".tmp.mp4"
        if os.path.exists(tmp_out):
            try:
                os.remove(tmp_out)
            except Exception:
                pass

        final_video.write_videofile(
            tmp_out,
            fps=TARGET_FPS,
            codec="libx264",
            audio_codec="aac",
            audio_fps=44100,
            preset="medium",
            verbose=False,
            logger=None,
            threads=2,
            ffmpeg_params=["-movflags", "+faststart", "-pix_fmt", "yuv420p"],
        )
        final_video.close()
        for c in video_clips:
            try:
                c.close()
            except Exception:
                pass

        if not os.path.exists(tmp_out):
            raise Exception("Video file was not created")

        # Prefer ffmpeg remux for reliable faststart when available
        if not _ensure_mp4_faststart(tmp_out, output_path):
            # Fallback: move temp file as-is (already requested +faststart above)
            if os.path.exists(output_path):
                try:
                    os.remove(output_path)
                except Exception:
                    pass
            os.replace(tmp_out, output_path)
        else:
            try:
                if os.path.exists(tmp_out):
                    os.remove(tmp_out)
            except Exception:
                pass

        if not os.path.exists(output_path):
            raise Exception("Video file was not created")

        file_size = os.path.getsize(output_path)
        result["success"] = True
        result["video_path"] = output_path
        result["duration"] = video_duration
        result["file_size"] = file_size
        provider_note = result["tts_provider"] or "unknown"
        if "silent" in provider_note and "elevenlabs" not in provider_note:
            result["message"] = (
                f"Video generated ({video_duration:.0f}s). "
                f"Voice fallback used ({provider_note}) — playable without ElevenLabs credits."
            )
        else:
            result["message"] = (
                f"Video generated successfully ({video_duration:.0f}s, "
                f"{file_size/1024/1024:.2f}MB, tts={provider_note})"
            )
        print(f"✓ Video saved: {output_path}")

    except Exception as e:
        result["error"] = str(e)
        result["message"] = f"Error generating video: {str(e)}"
        print(f"✗ Error: {str(e)}")
        import traceback
        traceback.print_exc()

    finally:
        for temp_dir in temp_dirs:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)

    return result


def get_or_create_video(lesson_id, lesson_notes_path, pptx_path, output_filename):
    ensure_video_directory()
    output_path = os.path.join(VIDEO_DIR, output_filename)

    if os.path.exists(output_path):
        file_size = os.path.getsize(output_path)
        return {
            "success": True,
            "video_path": output_path,
            "cached": True,
            "file_size": file_size,
            "message": "Video loaded from cache",
        }

    return generate_video(lesson_notes_path, pptx_path, output_path)


def list_generated_videos():
    ensure_video_directory()
    videos = []
    for filename in os.listdir(VIDEO_DIR):
        if filename.endswith(".mp4"):
            filepath = os.path.join(VIDEO_DIR, filename)
            file_size = os.path.getsize(filepath)
            mtime = os.path.getmtime(filepath)
            videos.append({
                "filename": filename,
                "filepath": filepath,
                "file_size": file_size,
                "created_at": datetime.fromtimestamp(mtime).isoformat(),
            })
    return sorted(videos, key=lambda x: x["created_at"], reverse=True)


def delete_video(video_filename):
    ensure_video_directory()
    filepath = os.path.join(VIDEO_DIR, video_filename)
    if os.path.exists(filepath) and filepath.startswith(VIDEO_DIR):
        try:
            os.remove(filepath)
            return True
        except Exception as e:
            print(f"Error deleting video: {e}")
            return False
    return False
