

from flask import Blueprint, abort, jsonify, render_template, request, redirect, url_for, flash, session, send_file, Response
import mysql
import os
import time
from werkzeug.security import generate_password_hash, check_password_hash
from app.utils.message import message
import json
import hashlib
from .db_management.db import get_db_connection
from .db_management.sql import insert as db_insert
from .db_management.sql import select as db_select
from .db_management.sql import update as db_update
from .db_management.sql import delete as db_delete
from .db_management.sql import select_one as db_select_one
from .services.methods import explain_text_service,summarize_service,text_to_speech_service,get_current_user_profile, save_course_reaction, is_enrolled, get_user_reaction, get_all_courses, get_recommended_courses, get_trending_courses, log_search, search_courses
from app.services.profile_service import *
# LLM API token and URL for LLM service
from app.services.llm_service import generate_course, generate_lesson_notes
import app.services.assessment_service as assessment_service
from app.services.notes.notes_service import generate_lesson_notes as create_lesson_notes, get_all_notes_for_course, download_notes
import traceback
from app.services.pptx.pptx_controller import generate_lesson_pptx
from app.services.pptx.pptx_service import sanitize_filename, PPTX_DIR
from app.services.video.video_controller import generate_lesson_video
from app.services.video.video_service import VIDEO_DIR
from app.services.methods import get_lesson_quiz_score, get_module_quiz_average, is_module_completed



main = Blueprint('main', __name__)

@main.route('/')
def index():
    return render_template('index.html')



@main.route('/login', methods=['GET', 'POST'])
def login():

    if request.method == 'GET':
        next_page = request.args.get("next")
        return render_template('login.html', next_page=next_page)

    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']
        next_page = request.form.get("next_page")  # 👈 IMPORTANT

        conn = get_db_connection()

        try:
            user = db_select_one(conn, """
                SELECT id, full_name, email, password_hash
                FROM users
                WHERE email = %s
            """, (email,))

            if user and check_password_hash(user[3], password):

                session['user_id'] = user[0]
                session['full_name'] = user[1]
                session['email'] = user[2]

                profile = get_current_user_profile()
                print("Profile object:", profile)
                print("Profile type:", type(profile))
                session["profile_level"] = profile["static_profile"]["explicit_level"]

                # Check if user has COMPLETE preferences (essential fields filled)
                pref = get_user_preferences(user[0], check_complete=True)
                print(f"User preferences (complete check): {pref}")

                if not pref:
                    # Preferences missing or incomplete - redirect to set them
                    print("Redirecting to set_preferences - no complete preferences found")
                    flash("Please complete your learning preferences to continue.", "info")
                    return redirect("/set_preferences")

                # Preferences exist and are complete - go to dashboard
                print("Complete preferences found - redirecting to dashboard")
                flash("Login successful!", "success")

                if next_page and next_page not in ("None", "", None):
                    return redirect(next_page)

                return redirect(url_for("main.student_dashboard"))
            else:
                flash("Invalid email or password.", "danger")

        except Exception as e:
            print (f" Error from db {str(e)}")
            print("Error from db:")
            traceback.print_exc()
            flash("Database error occurred.", "danger")

        finally:
            conn.close()

    return render_template('login.html')

@main.route('/register', methods=['GET', 'POST'])
def register():

    if request.method == 'POST':

        full_name = request.form['full_name']
        email = request.form['email']
        password = request.form['password']

        password_hash = generate_password_hash(password)

        conn = get_db_connection()
        cursor = conn.cursor()

        try:
            
            db_insert(conn, """
                INSERT INTO users (full_name, email, password_hash)
                VALUES (%s, %s, %s)
            """, (full_name, email, password_hash))
            conn.commit()

            message("Account created successfully!", "success")
            return redirect(url_for('main.login'))

        except mysql.connector.Error as err:
            message("Email already exists.", "danger")

        finally:
            cursor.close()
            conn.close()

    return render_template('register.html')

@main.route('/student_dashboard')
def student_dashboard():
    if 'user_id' not in session:
        message("Please log in to access the dashboard.", "warning")
        return redirect(url_for('main.login'))

    return render_template('/student/student_dashboard.html', full_name=session.get('full_name'))

@main.route('/preferences')
def preferences():
    return render_template('preference.html')

@main.route('/generate_preview', methods=['POST', 'GET'])
def generate_preview():

    preferences = {
        "domain": request.form.get("domain"),
        "topic": request.form.get("topic"),
        "goal": request.form.get("goal"),
        "level": request.form.get("level"),
        "duration": request.form.get("duration"),
        "learning_preference": request.form.get("learning_preference"),
        "prior_knowledge": request.form.get("prior_knowledge")
    }

    try:
        if not preferences["domain"] or not preferences["topic"]:
            flash("Domain and Topic are required fields.", "danger")
            return redirect(url_for("main.preferences"))

        course_data = generate_course(preferences)
        session["generated_course"] = course_data
        session["preferences"] = preferences
        print("Generated Course Data:", course_data)
        print (f'course in session: {session["generated_course"]} and preferences: {session["preferences"]}')
        

        if not course_data:
            flash("Course generation failed. Please try again.", "danger")
            return redirect(url_for("main.preferences"))

    except Exception as e:
        print("Error during course generation:", str(e))
        flash("An error occurred during course generation.", "danger")
        return redirect(url_for("main.preferences"))

    return render_template("preview.html", course=course_data)

@main.route('/loader_test')
def loader_test():
    """Test page for loader animation"""
    if "user_id" not in session:
        return redirect(url_for("main.login"))
    
    return render_template("student/loader_test.html")

@main.route('/preview_generated_course')
def preview_generated_course():
    """Show the generated course preview (keeps user in the flow)"""
    
    if "user_id" not in session:
        flash("Please log in to view your generated course.", "warning")
        return redirect(url_for("main.login"))
    
    course_data = session.get("generated_course")
    
    if not course_data:
        flash("No generated course found. Please generate a course first.", "warning")
        return redirect(url_for("main.discover"))
    
    # Check if user has already reacted to this course
    conn = get_db_connection()
    try:
        content_hash = hashlib.sha256(
            json.dumps(course_data, sort_keys=True).encode("utf-8")
        ).hexdigest()
        
        course = db_select_one(conn, """
            SELECT id
            FROM courses
            WHERE content_hash = %s
        """, (content_hash,))
        
        user_reaction = None
        is_enrolled = False
        course_id_from_db = None
        
        if course:
            course_id_from_db = course[0]
            
            # Check reaction
            feedback = db_select_one(conn, """
                SELECT reaction
                FROM course_feedback
                WHERE user_id = %s AND course_id = %s
            """, (session["user_id"], course_id_from_db))
            
            if feedback:
                user_reaction = feedback[0]
            
            # Check enrollment
            enrollment = db_select_one(conn, """
                SELECT id
                FROM enrollments
                WHERE student_id = %s AND course_id = %s
            """, (session["user_id"], course_id_from_db))
            
            if enrollment:
                is_enrolled = True
        
        return render_template("preview.html", 
                             course=course_data, 
                             user_reaction=user_reaction,
                             is_enrolled=is_enrolled,
                             course_id=course_id_from_db)
        
    finally:
        conn.close()

@main.route('/react_course', methods=['POST'])
def react_course():

    if "user_id" not in session:
        flash("Please log in to react to this course.", "warning")

        next_page = request.form.get("next_page")

        print("NEXT PAGE FROM FORM:", next_page)  # debug

        return redirect(url_for("main.login", next=next_page))

    # Save the reaction and get the response
    response = save_course_reaction()
    
    # If it's a redirect, follow it
    return response

@main.route('/enroll_generated_course')
def enroll_generated_course():
    """Enroll in the currently generated course (from session)"""
    
    if "user_id" not in session:
        flash("Please log in to enroll in this course.", "warning")
        return redirect(url_for("main.login"))
    
    course_data = session.get("generated_course")
    
    if not course_data:
        flash("No generated course found. Please generate a course first.", "warning")
        return redirect(url_for("main.preview_generated_course"))
    
    conn = get_db_connection()
    
    try:
        # Get course ID from content hash
        content_hash = hashlib.sha256(
            json.dumps(course_data, sort_keys=True).encode("utf-8")
        ).hexdigest()
        
        course = db_select_one(conn, """
            SELECT id
            FROM courses
            WHERE content_hash = %s
        """, (content_hash,))
        
        if not course:
            flash("Course not found in database. Please like/dislike first to save it.", "warning")
            return redirect(url_for("main.preview_generated_course"))
        
        course_id = course[0]
        
        # Enroll user
        db_insert(conn, """
            INSERT IGNORE INTO enrollments (user_id, course_id)
            VALUES (%s, %s)
        """, (session["user_id"], course_id))
        
        flash("🎉 Successfully enrolled! This course is now saved in your account.", "success")
        return redirect(url_for("main.my_courses"))
        
    except Exception as e:
        print(f"Error enrolling in generated course: {e}")
        traceback.print_exc()
        flash("An error occurred while enrolling.", "danger")
        return redirect(url_for("main.preview_generated_course"))
    
    finally:
        conn.close()


@main.route("/courses")
def courses():

    conn = get_db_connection()

    try:
        courses = db_select(conn, """
            SELECT id, title, description, popularity_score, created_at
            FROM courses
            WHERE is_public = 1
            ORDER BY popularity_score DESC, created_at DESC
        """)

        formatted_courses = []

        for course in courses:
            formatted_courses.append({
                "id": course[0],
                "title": course[1],
                "description": course[2],
                "popularity_score": course[3],
                "created_at": course[4]
            })

        # ✅ ADD THIS BLOCK HERE
        formatted_courses = []

        # Get enrolled IDs once
        if "user_id" in session:
            enrolled_courses = db_select(conn, """
                SELECT course_id FROM enrollments
                WHERE user_id = %s
            """, (session["user_id"],))

            enrolled_ids = {row[0] for row in enrolled_courses}
        else:
            enrolled_ids = set()

        for course in courses:
            course_id = course[0]

            reaction = get_user_reaction(session.get("user_id"), course_id) if "user_id" in session else None
            enrolled = course_id in enrolled_ids

            formatted_courses.append({
                "id": course_id,
                "title": course[1],
                "description": course[2],
                "popularity_score": course[3],
                "created_at": course[4],
                "reaction": reaction,
                "enrolled": enrolled
            })
      

        return render_template(
            "courses.html",
            courses=formatted_courses,
            
        )

    finally:
        conn.close()

@main.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('main.index'))


@main.route("/course/<int:course_id>")
def view_course(course_id):
    
    if 'user_id' not in session:
        # Store the intended destination and redirect to login
        return redirect(url_for("main.login", next=url_for("main.view_course", course_id=course_id)))

    conn = get_db_connection()

    try:
        # 1️⃣ Fetch course
        course = db_select_one(conn, """
            SELECT id, title, description, content, popularity_score, created_at
            FROM courses
            WHERE id = %s
        """, (course_id,))

        if not course:
            abort(404)

        # 2️⃣ Unpack
        course_data = {
            "id": course[0],
            "title": course[1],
            "description": course[2],
            "content": json.loads(course[3]),
            "popularity_score": course[4],
            "created_at": course[5]
        }

        # 3️⃣ Check if logged-in user reacted
        user_reaction = None

        feedback = db_select_one(conn, """
            SELECT reaction
            FROM course_feedback
            WHERE user_id = %s AND course_id = %s
        """, (session["user_id"], course_id))

        if feedback:
            user_reaction = feedback[0]

        # 4️⃣ Check enrollment status
        enrolled = is_enrolled(session["user_id"], course_id)

        return render_template(
            "course_content.html",
            course=course_data,
            user_reaction=user_reaction,
            enrolled=enrolled
        )

    finally:
        conn.close()

# enrollment route
@main.route("/enroll/<int:course_id>")
def enroll(course_id):

    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.enroll", course_id=course_id)))

    conn = get_db_connection()

    try:
        db_insert(conn, """
            INSERT IGNORE INTO enrollments (user_id, course_id)
            VALUES (%s, %s)
        """, (session["user_id"], course_id))

        flash("Successfully enrolled in course.", "success")

        return redirect(url_for("main.view_course", course_id=course_id))

    finally:
        conn.close()


@main.route('/generate_lesson_notes/<int:course_id>/<int:module_index>/<int:lesson_index>')
def generate_lesson_notes_route(course_id, module_index, lesson_index):
    """Generate notes for a specific lesson"""
    
    if "user_id" not in session:
        flash("Please log in to generate lesson notes.", "warning")
        return redirect(url_for("main.login"))
    
    conn = get_db_connection()
    
    try:
        # Fetch course from database
        course = db_select_one(conn, """
            SELECT id, title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))
        
        if not course:
            flash("Course not found.", "danger")
            return redirect(url_for("main.my_courses"))
        
        course_title = course[1]
        course_content = json.loads(course[2])
        
        # Get user preferences for personalization
        user_preferences = get_user_preferences(session["user_id"])
        
        # Extract module and lesson
        modules = course_content.get('modules', [])
        
        if module_index < 0 or module_index >= len(modules):
            flash("Invalid module.", "danger")
            return redirect(url_for("main.my_courses"))
        
        module = modules[module_index]
        module_title = module.get('title', f'Module {module_index + 1}')
        
        lessons = module.get('lessons', [])
        
        if lesson_index < 0 or lesson_index >= len(lessons):
            flash("Invalid lesson.", "danger")
            return redirect(url_for("main.my_courses"))
        
        lesson = lessons[lesson_index]
        lesson_title = lesson.get('title', f'Lesson {lesson_index + 1}')
        
        # CHECK IF NOTES ALREADY EXIST
        from app.services.notes.notes_service import get_existing_notes_file
        existing_file = get_existing_notes_file(course_title, module_title, lesson_title)
        
        if existing_file and os.path.exists(existing_file):
            print(f"Notes already exist: {existing_file}")
            flash("Using existing notes.", "info")
            
            # Store existing file path in session
            session[f'notes_{course_id}_{module_index}_{lesson_index}'] = existing_file
            
            # Redirect to view/download options
            flash("Notes are ready! You can view online or download.", "success")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        
        # Notes don't exist - generate new ones
        print(f"Generating NEW notes for: {course_title} > {module_title} > {lesson_title}")
        
        # Generate notes using LLM
        notes_result = create_lesson_notes(
            course_title=course_title,
            module_title=module_title,
            lesson=lesson,  # Pass the entire lesson dict
            preferences=user_preferences
        )
        
        if notes_result and notes_result.get('success'):
            flash("Lesson notes generated successfully!", "success")
            
            # Store file path in session for both view and download
            session[f'notes_{course_id}_{module_index}_{lesson_index}'] = notes_result['file_path']
            
            # Show success with options - redirect to a choice page or back to course
            # For now, let's redirect back to course with a message
            flash("You can now view online or download the notes.", "info")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        else:
            error_msg = notes_result.get('message', 'Failed to generate notes') if notes_result else 'Unknown error'
            flash(f"Error generating notes: {error_msg}", "danger")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
            
    except Exception as e:
        print(f"Error generating lesson notes: {str(e)}")
        traceback.print_exc()
        flash("An error occurred while generating notes.", "danger")
        return redirect(url_for("main.my_courses"))
    
    finally:
        conn.close()

@main.route ('/view_lesson_slides/<int:course_id>/<int:module_index>/<int:lesson_index>')
def view_lesson_slides(course_id, module_index, lesson_index):
    """View generated lesson slides in browser"""
    
    if "user_id" not in session:
        flash("Please log in to view lesson slides.", "warning")
        return redirect(url_for("main.login"))
    
    try:
        # Get file path from session (note: stored as 'pptx_*' not 'slides_*')
        file_path = session.get(f'pptx_{course_id}_{module_index}_{lesson_index}')
        
        if not file_path or not os.path.exists(file_path):
            flash("Slides not generated yet. Please generate slides first.", "warning")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        
        # Get filename from path
        filename = os.path.basename(file_path)
        title = filename.replace('.pptx', '').replace('_', ' ').title()
        
        # Render viewer template with PPTX file path
        return render_template("student/lesson_slides_viewer.html", 
                             file_path=file_path,
                             filename=filename,
                             title=title,
                             course_id=course_id,
                             module_index=module_index,
                             lesson_index=lesson_index)
            
    except Exception as e:
        print(f"Error viewing slides: {str(e)}")
        flash("Error loading slides.", "danger")
        return redirect(url_for("main.my_courses"))
    
@main.route('/view_lesson_notes/<int:course_id>/<int:module_index>/<int:lesson_index>')
def view_lesson_notes(course_id, module_index, lesson_index):
    """View lesson notes in browser (without downloading)"""
    
    if "user_id" not in session:
        flash("Please log in to view lesson notes.", "warning")
        return redirect(url_for("main.login"))
    
    try:
        # Get file path from session
        file_path = session.get(f'notes_{course_id}_{module_index}_{lesson_index}')
        
        if not file_path or not os.path.exists(file_path):
            flash("Notes not generated yet. Please generate notes first.", "warning")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        
        # Read the content from the Word document
        from docx import Document
        doc = Document(file_path)
        
        # Extract all paragraphs
        content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Safely check formatting
                is_bold = False
                is_italic = False
                is_heading = False
                style_name_str = ''
                
                # Check runs for bold/italic
                try:
                    for run in paragraph.runs:
                        if hasattr(run, 'bold') and run.bold is True:
                            is_bold = True
                        if hasattr(run, 'italic') and run.italic is True:
                            is_italic = True
                except Exception as e:
                    print(f"Error checking run formatting: {e}")
                
                # Check if it's a heading style
                try:
                    if hasattr(paragraph, 'style') and hasattr(paragraph.style, 'name'):
                        style_name_obj = paragraph.style.name
                        if style_name_obj is not None and isinstance(style_name_obj, str):
                            style_name_str = style_name_obj
                            is_heading = style_name_str.startswith('Heading')
                except Exception as e:
                    print(f"Error checking style name: {e}")
                
                content.append({
                    'text': paragraph.text,
                    'bold': is_bold,
                    'italic': is_italic,
                    'heading': is_heading,
                    'style_name': style_name_str
                })
        
        # Get document title from first heading or filename
        title = os.path.basename(file_path).replace('.docx', '').replace('_', ' ').title()
        
        return render_template("student/lesson_notes_viewer.html", 
                             content=content, 
                             title=title,
                             course_id=course_id,
                             module_index=module_index,
                             lesson_index=lesson_index)
            
    except Exception as e:
        print(f"Error viewing notes: {str(e)}")
        flash("Error loading notes.", "danger")
        return redirect(url_for("main.my_courses"))


@main.route('/download_lesson_notes/<int:course_id>/<int:module_index>/<int:lesson_index>')
def download_lesson_notes(course_id, module_index, lesson_index):
    """Download generated lesson notes"""
    
    if "user_id" not in session:
        flash("Please log in to download notes.", "warning")
        return redirect(url_for("main.login"))
    
    try:
        # Get file path from session
        file_path = session.get(f'notes_{course_id}_{module_index}_{lesson_index}')
        
        if not file_path or not os.path.exists(file_path):
            flash("Notes file not found. Please generate notes first.", "warning")
            return redirect(url_for("main.my_courses"))
        
        # Read file
        filename, content = download_notes(file_path)
        
        if filename and content:
            from flask import send_file
            return send_file(
                file_path,
                as_attachment=True,
                download_name=filename,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
        else:
            flash("Error reading notes file.", "danger")
            return redirect(url_for("main.my_courses"))
            
    except Exception as e:
        print(f"Error downloading notes: {str(e)}")
        flash("Error downloading notes.", "danger")
        return redirect(url_for("main.my_courses"))


@main.route('/lesson_quiz')
def lesson_quiz():
    """Render the lesson quiz page using query parameters."""
    if "user_id" not in session:
        flash("Please log in to take the quiz.", "warning")
        return redirect(url_for("main.login"))

    course_id = request.args.get('course')
    module_index = request.args.get('module')
    lesson_index = request.args.get('lesson')

    if course_id is None or module_index is None or lesson_index is None:
        flash("Invalid quiz request. Please select a lesson quiz.", "danger")
        return redirect(url_for("main.my_courses"))

    try:
        course_id = int(course_id)
        module_index = int(module_index)
        lesson_index = int(lesson_index)
    except ValueError:
        flash("Invalid quiz identifiers.", "danger")
        return redirect(url_for("main.my_courses"))


    
    return render_template(
        'student/lesson_quiz.html',
        course_id=course_id,
        module_index=module_index,
        lesson_index=lesson_index
    )


@main.route('/generate_lesson_quiz/<int:course_id>/<int:module_index>/<int:lesson_index>')
def generate_lesson_quiz(course_id, module_index, lesson_index):
    """Generate quiz for a specific lesson"""
    
    if "user_id" not in session:
        flash("Please log in to generate quiz.", "warning")
        return redirect(url_for("main.login"))
    
    conn = get_db_connection()
    
    try:
        # Fetch course from database
        course = db_select_one(conn, """
            SELECT id, title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))
        
        if not course:
            flash("Course not found.", "danger")
            return redirect(url_for("main.my_courses"))
        
        course_title = course[1]
        course_content = json.loads(course[2])
        
        # Get lesson data
        try:
            lesson_data = course_content['modules'][module_index]['lessons'][lesson_index]
        except (IndexError, KeyError):
            flash("Lesson not found.", "danger")
            return redirect(url_for("main.view_my_course", course_id=course_id))
        
        # Generate quiz
        try:
            question_ids = assessment_service.create_quiz_for_lesson(course_id, module_index, lesson_index, lesson_data)
            
            if question_ids and len(question_ids) > 0:
                flash(f"Quiz generated successfully with {len(question_ids)} questions!", "success")
            else:
                flash("Quiz generation completed but no questions were created.", "warning")
                
        except Exception as e:
            print(f"Error generating quiz: {str(e)}")
            traceback.print_exc()
            flash("Error generating quiz. Please try again.", "danger")
        
        return redirect(url_for("main.view_my_course", course_id=course_id))
            
    except Exception as e:
        print(f"Error in quiz generation: {str(e)}")
        traceback.print_exc()
        flash("An error occurred while generating quiz.", "danger")
        return redirect(url_for("main.my_courses"))
    
    finally:
        conn.close()


@main.route("/my_courses")
def my_courses():

    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.my_courses")))

    conn = get_db_connection()

    try:
        courses = db_select(conn, """
            SELECT c.id, c.title, c.content, c.popularity_score, c.created_at
            FROM courses c
            JOIN enrollments e ON c.id = e.course_id
            WHERE e.user_id = %s
            ORDER BY e.enrolled_at DESC
        """, (session["user_id"],))

        formatted_courses = []

        for course in courses:
            formatted_courses.append({
                "id": course[0],
                "title": course[1],
                "overview": json.loads(course[2]).get("overview", ""),
                "popularity_score": course[3],
                "created_at": course[4]
            })

        return render_template("student/courses.html", courses=formatted_courses)

    finally:
        conn.close()


@main.route("/view_my_course/<int:course_id>")
def view_my_course(course_id):

    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.view_my_course", course_id=course_id)))

    conn = get_db_connection()

    try:
        course = db_select_one(conn, """
            SELECT id, title, description, content, popularity_score, created_at
            FROM courses
            WHERE id = %s
        """, (course_id,))

        if not course:
            abort(404)

        course_data = {
            "id": course[0],
            "title": course[1],
            "description": course[2],
            "content": json.loads(course[3]),
            "popularity_score": course[4],
            "created_at": course[5]
        }
        reaction = get_user_reaction(session.get("user_id"), course_data["id"]) if "user_id" in session else None
        print("User reaction for course:", reaction)  # debug
        enrolled = is_enrolled(session.get("user_id"), course_data["id"]) if "user_id" in session else False
        
        # Build scores dict for each lesson
        scores = {}
        course_content = course_data["content"]
        for module_idx, module in enumerate(course_content.get('modules', [])):
            for lesson_idx, lesson in enumerate(module.get('lessons', [])):
                scores[f"{module_idx}_{lesson_idx}"] = get_lesson_quiz_score(session["user_id"], course_id, module_idx, lesson_idx)
        
        # Also get module averages and completion status
        module_averages = {}
        module_completed = {}
        for module_idx in range(len(course_content.get('modules', []))):
            module_averages[module_idx] = get_module_quiz_average(session["user_id"], course_id, module_idx)
            module_completed[module_idx] = is_module_completed(session["user_id"], course_id, module_idx)
        
        return render_template("student/view_course.html", scores=scores, module_averages=module_averages, module_completed=module_completed, course=course_data, reaction=reaction, enrolled=enrolled)

    finally:
        conn.close()

@main.route('/module_assessment/<int:course_id>/<int:module_index>')
def module_assessment(course_id, module_index):
    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.module_assessment", course_id=course_id, module_index=module_index)))
    
    # Check if user is enrolled
    if not is_enrolled(session["user_id"], course_id):
        flash("You must be enrolled in this course to access assessments.", "error")
        return redirect(url_for("main.view_my_course", course_id=course_id))
    
    # Check if module is completed (all lessons have quiz results)
    if not is_module_completed(session["user_id"], course_id, module_index):
        flash("You must complete all lesson quizzes in this module before taking the module assessment.", "warning")
        return redirect(url_for("main.view_my_course", course_id=course_id))
    
    conn = get_db_connection()
    try:
        course = db_select_one(conn, """
            SELECT title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))
        
        if not course:
            abort(404)
        
        course_title = course[0]
        course_content = json.loads(course[1])
        
        # Get module title
        module_title = f"Module {module_index + 1}"
        try:
            module = course_content['modules'][module_index]
            module_title = module.get('title', module_title)
        except (IndexError, KeyError):
            pass
        
        return render_template("student/module_assessment.html", 
                             course_id=course_id, 
                             course_title=course_title,
                             module_index=module_index, 
                             module_title=module_title)
    
    finally:
        conn.close()

@main.route('/enroll_course/<int:course_id>', methods=['POST', 'GET'])
def enroll_course(course_id):
    if "user_id" not in session:
        return redirect(url_for("main.login", next=url_for("main.enroll_course", course_id=course_id)))

    conn = get_db_connection()

    try:
        db_insert(conn, """
            INSERT IGNORE INTO enrollments (user_id, course_id)
            VALUES (%s, %s)
        """, (session["user_id"], course_id))

        flash("Successfully enrolled in course.", "success")

        return redirect(url_for("main.view_my_course", course_id=course_id))

    finally:
        conn.close()


@main.route('/discover')
def discover():

    if 'user_id' not in session:
        return redirect(url_for("main.login"))
    student_id = session['user_id']
    query = request.args.get('query', '').strip()

    existing_courses = get_all_courses()
    search_results = []
    recommended_courses = get_recommended_courses(student_id)
    trending_courses = get_trending_courses()
    

    if query:
        # Log search behavior
        log_search(student_id, query)

        # Intelligent search
        search_results = search_courses(query)
        reaction = get_user_reaction(session.get("user_id"), trending_courses[0]["id"]) if "user_id" in session and trending_courses else None
        print("User reaction for trending course:", reaction)  # debug
        enrolled = is_enrolled(session.get("user_id"), trending_courses[0]["id"]) if "user_id" in session and trending_courses else False

    return render_template(
        "student/discover.html",
        search_results=search_results,
        recommended_courses=recommended_courses,
        trending_courses=trending_courses,
        existing_courses=existing_courses,
        query=query, 
        reaction=reaction if query else None,
        enrolled=enrolled if query else None
    )
# discover course as a student where the course they are looking for is not available yet, so they can input their preferences and get a generated course preview. Note this, we are using search behavior to trigger the course generation flow, so we will have a search bar on the discover page where they can input their desired course topic, and if it does not exist in the database, we will redirect them to the preferences page with the topic pre-filled. 
@main.route('/search', methods=['POST'])
def search():

    topic = request.form.get("topic")

    if not topic:
        flash("Please enter a topic to search.", "warning")
        return redirect(url_for("main.discover"))

    conn = get_db_connection()

    try:
        course = db_select_one(conn, """
            SELECT id
            FROM courses
            WHERE JSON_EXTRACT(content, '$.topic') = %s
        """, (json.dumps(topic),))

        if course:
            return redirect(url_for("main.view_course", course_id=course[0]))
        else:
            flash("Course not found. Please provide your preferences to generate a course preview.", "info")
            return redirect(url_for("main.set_preferences", topic=topic))

    finally:
        conn.close()


# Profile Building 
@main.route("/set_preferences", methods=["GET", "POST"])
def set_preferences():

    if "user_id" not in session:
        return redirect("/login")

    user_id = session["user_id"]

    if request.method == "POST":

        data = {
            "domain": request.form.get("domain"),
            "topic": request.form.get("topic"),
            "goal": request.form.get("goal"),
            "level": request.form.get("level"),
            "duration": request.form.get("duration"),
            "learning_preference": request.form.get("learning_preference"),
            "prior_knowledge": request.form.get("prior_knowledge")
        }

        save_user_preferences(user_id, data)

        flash("Your preferences saved successfully", "success")
        return redirect(url_for("main.student_dashboard"))

    existing = get_user_preferences(user_id)

    return render_template(
        "student/preferences.html",
        existing=existing
    )
#generate course from logged in student 
@main.route ("/learner_generate_course", methods=["GET", "POST"])
def learner_generate_course(): 
    if "user_id" not in session:
        return redirect(url_for("main.login"))
    
    user_id = session["user_id"]
    topic = request.args.get("topic") or request.form.get("topic")
    
    conn = get_db_connection()
    
    try:
        # Get user preferences
        user_preferences = get_user_preferences(user_id)
        
        if not user_preferences:
            flash("Please set your learning preferences first.", "warning")
            return redirect(url_for("main.set_preferences"))
        
        # If POST request, generate the course
        if request.method == "POST":
            # Build comprehensive preferences including the search query
            generation_prefs = {
                "domain": user_preferences.get("domain", ["General"]),
                "topic": topic or user_preferences.get("topic", ""),
                "goal": user_preferences.get("goal", "Professional Skill Development"),
                "level": user_preferences.get("level", "Beginner"),
                "duration": user_preferences.get("duration", "4"),
                "learning_preference": user_preferences.get("learning_preference", "Balanced Approach"),
                "prior_knowledge": user_preferences.get("prior_knowledge", "")
            }
            
            print(f"Generating course with preferences: {generation_prefs}")
            
            # Generate course using LLM
            course_data = generate_course(generation_prefs)
            
            if course_data:
                session["generated_course"] = course_data
                session["preferences"] = generation_prefs
                
                # Log the search behavior for personalization
                log_search(user_id, topic or generation_prefs["topic"])
                
                # Auto-enroll student in the generated course
                try:
                    content_hash = hashlib.sha256(
                        json.dumps(course_data, sort_keys=True).encode("utf-8")
                    ).hexdigest()
                    
                    # Check if course exists in database
                    existing_course = db_select_one(conn, """
                        SELECT id
                        FROM courses
                        WHERE content_hash = %s
                    """, (content_hash,))
                    
                    if existing_course:
                        # Course exists - enroll student
                        course_id = existing_course[0]
                        
                        # Check if already enrolled
                        enrollment = db_select_one(conn, """
                            SELECT id
                            FROM enrollments
                            WHERE student_id = %s AND course_id = %s
                        """, (user_id, course_id))
                        
                        if not enrollment:
                            db_insert(conn, """
                                INSERT INTO enrollments (student_id, course_id)
                                VALUES (%s, %s)
                            """, (user_id, course_id))
                    
                except Exception as e:
                    print(f"Error auto-enrolling: {e}")
                    # Continue anyway - enrollment can happen later
                
                flash("Course generated successfully!", "success")
                return render_template("preview.html", course=course_data)
            else:
                flash("Failed to generate course. Please try again.", "danger")
                
    except Exception as err:
        print(f"Error generating course: {err}")
        traceback.print_exc()
        flash("An error occurred during course generation.", "danger")
    
    finally:
        conn.close()
    
    # GET request - show the generation page
    return render_template(
        "student/personalised_course.html", 
        user_preferences=user_preferences,
        topic=topic
    )

@main.route('/user_profile')
def user_profile():
    profile_data = build_user_profile(session["user_id"])
    return render_template("student/profile.html", profile=profile_data)

@main.route('/text-to-speech', methods=['POST'])
def text_to_speech():
    try:
        data = request.get_json()
        text = data.get('text')
        
        if not text:
            return jsonify({"error": "No text provided"}), 400

        audio = text_to_speech_service(text)

        # Return the raw bytes as an MPEG stream
        return Response(audio, mimetype='audio/mpeg')

    except Exception as e:
        print(f"CRITICAL TTS ERROR: {str(e)}") # This shows in your terminal
        return jsonify({"error": str(e)}), 500
    

@main.route('/summarize', methods=['POST'])
def summarize():
    try:
        data = request.get_json()
        text = data.get('text')

        summary = summarize_service(text)

        return jsonify({"summary": summary})

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    
@main.route('/explain', methods=['POST'])
def explain():
    try:
        data = request.get_json()
        text = data.get('text')

        explanation = explain_text_service(text)

        return jsonify({"explanation": explanation})

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    
# Route to generate PPTX for a specific lesson
@main.route('/generate_lesson_pptx/<int:course_id>/<int:module_index>/<int:lesson_index>')
def generate_lesson_pptx_route(course_id, module_index, lesson_index):
    """Generate PPTX for a specific lesson"""

    if "user_id" not in session:
        flash("Please log in to generate presentation.", "warning")
        return redirect(url_for("main.login"))

    conn = get_db_connection()

    try:
        # Fetch course
        course = db_select_one(conn, """
            SELECT id, title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))

        if not course:
            flash("Course not found.", "danger")
            return redirect(url_for("main.my_courses"))

        course_title = course[1]
        course_content = json.loads(course[2])

        # User preferences
        user_preferences = get_user_preferences(session["user_id"])

        # Extract module
        modules = course_content.get('modules', [])
        if module_index < 0 or module_index >= len(modules):
            flash("Invalid module.", "danger")
            return redirect(url_for("main.my_courses"))

        module = modules[module_index]
        module_title = module.get('title', f'Module {module_index + 1}')

        # Extract lesson
        lessons = module.get('lessons', [])
        if lesson_index < 0 or lesson_index >= len(lessons):
            flash("Invalid lesson.", "danger")
            return redirect(url_for("main.my_courses"))

        lesson = lessons[lesson_index]
        lesson_title = lesson.get('title', f'Lesson {lesson_index + 1}')

        # 🔷 CHECK IF PPTX ALREADY EXISTS
        
        filename = f"{sanitize_filename(course_title)}_{sanitize_filename(module_title)}_{sanitize_filename(lesson_title)}.pptx"
        file_path = os.path.join(PPTX_DIR, filename)

        if os.path.exists(file_path):
            print(f"PPTX already exists: {file_path}")
            flash("Using existing presentation.", "info")

            session[f'pptx_{course_id}_{module_index}_{lesson_index}'] = file_path

            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))

        # 🔷 GENERATE NEW PPTX
        print(f"Generating PPTX for: {course_title} > {module_title} > {lesson_title}")


        pptx_result = generate_lesson_pptx(
            course_title=course_title,
            module_title=module_title,
            lesson=lesson,
            preferences=user_preferences
        )

        if pptx_result and pptx_result.get('success'):
            flash("Presentation generated successfully!", "success")

            session[f'pptx_{course_id}_{module_index}_{lesson_index}'] = pptx_result['file_path']

            flash("You can now view or download the presentation.", "info")

            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))
        else:
            error_msg = pptx_result.get('message', 'Failed to generate PPTX') if pptx_result else 'Unknown error'
            flash(f"Error generating PPTX: {error_msg}", "danger")

            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))

    except Exception as e:
        print(f"Error generating PPTX: {str(e)}")
        traceback.print_exc()

        flash("An error occurred while generating presentation.", "danger")
        return redirect(url_for("main.my_courses"))

    finally:
        conn.close()

@main.route('/download_lesson_pptx/<int:course_id>/<int:module_index>/<int:lesson_index>')
def download_lesson_pptx(course_id, module_index, lesson_index):
    """Download generated PPTX file"""

    if "user_id" not in session:
        flash("Please log in to download slides.", "warning")
        return redirect(url_for("main.login"))

    try:
        key = f'pptx_{course_id}_{module_index}_{lesson_index}'
        file_path = session.get(key)

        if file_path and os.path.exists(file_path):
            return send_file(
                file_path,
                as_attachment=True,
                download_name=os.path.basename(file_path),
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        flash("Presentation not found. Please generate it first.", "danger")
        return redirect(url_for(
            'main.view_my_course',
            course_id=course_id
        ))

    except Exception as e:
        print(f"Error downloading PPTX: {str(e)}")
        traceback.print_exc()

        flash("An error occurred while downloading presentation.", "danger")
        return redirect(url_for("main.my_courses"))

@main.route('/serve_lesson_pptx/<int:course_id>/<int:module_index>/<int:lesson_index>')
def serve_lesson_pptx(course_id, module_index, lesson_index):
    """Serve PPTX slides as images for browser viewing"""
    
    if "user_id" not in session:
        flash("Please log in to view slides.", "warning")
        return redirect(url_for("main.login"))
    
    try:
        key = f'pptx_{course_id}_{module_index}_{lesson_index}'
        pptx_path = session.get(key)
        
        if not pptx_path or not os.path.exists(pptx_path):
            return jsonify({"error": "Presentation not found"}), 404
        
        # Convert PPTX to images
        slide_images = convert_pptx_to_images(pptx_path)
        
        if not slide_images:
            return jsonify({"error": "Could not convert presentation"}), 500
        
        return jsonify({
            "success": True,
            "slides": slide_images,
            "total_slides": len(slide_images)
        })
        
    except Exception as e:
        print(f"Error serving PPTX: {str(e)}")
        return jsonify({"error": str(e)}), 500

def convert_pptx_to_images(pptx_path):
    """Convert PPTX slides to PNG images"""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        import io
        import base64
        
        prs = Presentation(pptx_path)
        slide_images = []
        
        # Set slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Convert EMU (English Metric Units) to pixels (96 DPI)
        px_width = int(slide_width * 96 / 914400)
        px_height = int(slide_height * 96 / 914400)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            try:
                # Create blank image with white background
                img = Image.new('RGB', (px_width, px_height), color='white')
                draw = ImageDraw.Draw(img)
                
                # Try to use a default font, fall back to default if not available
                try:
                    font = ImageFont.truetype("arial.ttf", 24)
                    small_font = ImageFont.truetype("arial.ttf", 16)
                except:
                    font = ImageFont.load_default()
                    small_font = font
                
                y_pos = 40
                
                # Extract and draw text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Estimate if this is a title (usually first shape with text)
                        if y_pos == 40:
                            draw.text((40, y_pos), text, fill='black', font=font)
                            y_pos += 60
                        else:
                            # Wrap text for readability
                            words = text.split()
                            line = ""
                            for word in words:
                                test_line = line + word + " "
                                if len(test_line) > 50:
                                    draw.text((40, y_pos), line, fill='black', font=small_font)
                                    y_pos += 25
                                    line = word + " "
                                else:
                                    line = test_line
                            if line:
                                draw.text((40, y_pos), line, fill='black', font=small_font)
                                y_pos += 25
                
                # Convert to base64
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
                
                slide_images.append({
                    "slide_num": slide_num,
                    "data": f"data:image/png;base64,{img_base64}"
                })
                
            except Exception as e:
                print(f"Error converting slide {slide_num}: {e}")
                # Create error slide
                img = Image.new('RGB', (px_width, px_height), color='white')
                draw = ImageDraw.Draw(img)
                draw.text((40, 40), f"Error rendering slide {slide_num}", fill='red')
                
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
                
                slide_images.append({
                    "slide_num": slide_num,
                    "data": f"data:image/png;base64,{img_base64}"
                })
        
        return slide_images
        
    except ImportError:
        print("python-pptx not installed. Install with: pip install python-pptx pillow")
        return None
    except Exception as e:
        print(f"Error converting PPTX to images: {e}")
        return None


# ==========================================
# VIDEO GENERATION ROUTES
# ==========================================

@main.route('/generate_lesson_video/<int:course_id>/<int:module_index>/<int:lesson_index>')
def generate_lesson_video_route(course_id, module_index, lesson_index):
    """Generate video from lesson notes and slides"""

    if "user_id" not in session:
        flash("Please log in to generate videos.", "warning")
        return redirect(url_for("main.login"))

    try:
        # Get course data
        conn = get_db_connection()
        course = db_select_one(conn, """
            SELECT id, title, content
            FROM courses
            WHERE id = %s
        """, (course_id,))

        if not course:
            flash("Course not found.", "danger")
            return redirect(url_for("main.my_courses"))

        # Parse course content
        course_data = json.loads(course[2]) if course[2] else {}
        modules = course_data.get('modules', [])

        if module_index >= len(modules):
            flash("Module not found.", "danger")
            return redirect(url_for("main.my_courses"))

        module = modules[module_index]
        lessons = module.get('lessons', [])

        if lesson_index >= len(lessons):
            flash("Lesson not found.", "danger")
            return redirect(url_for("main.my_courses"))

        lesson = lessons[lesson_index]

        # Check if notes and PPTX exist
        notes_key = f'notes_{course_id}_{module_index}_{lesson_index}'
        pptx_key = f'pptx_{course_id}_{module_index}_{lesson_index}'

        notes_path = session.get(notes_key)
        pptx_path = session.get(pptx_key)

        if not notes_path or not os.path.exists(notes_path):
            flash("Lesson notes not found. Please generate notes first.", "warning")
            return redirect(url_for('main.view_my_course', course_id=course_id))

        if not pptx_path or not os.path.exists(pptx_path):
            flash("Presentation slides not found. Please generate slides first.", "warning")
            return redirect(url_for('main.view_my_course', course_id=course_id))

        # Generate video
        result = generate_lesson_video(
            course[1],  # course_title
            module['title'],
            lesson,
            notes_path,
            pptx_path
        )

        if result['success']:
            # Store video path in session
            video_key = f'video_{course_id}_{module_index}_{lesson_index}'
            session[video_key] = result['video_path']

            flash(f"Video generated successfully! Duration: {result['duration']:.0f}s", "success")
        else:
            flash(f"Error generating video: {result['message']}", "danger")

        return redirect(url_for('main.view_my_course', course_id=course_id))

    except Exception as e:
        print(f"Error generating video: {str(e)}")
        traceback.print_exc()
        flash("An error occurred while generating video.", "danger")
        return redirect(url_for("main.my_courses"))

    finally:
        conn.close()


@main.route('/view_lesson_video/<int:course_id>/<int:module_index>/<int:lesson_index>')
def view_lesson_video(course_id, module_index, lesson_index):
    """View generated lesson video in browser"""

    if "user_id" not in session:
        flash("Please log in to view lesson videos.", "warning")
        return redirect(url_for("main.login"))

    try:
        # Get file path from session
        file_path = session.get(f'video_{course_id}_{module_index}_{lesson_index}')

        if not file_path or not os.path.exists(file_path):
            flash("Video not generated yet. Please generate video first.", "warning")
            return redirect(url_for(
                'main.view_my_course',
                course_id=course_id
            ))

        # Get filename from path
        filename = os.path.basename(file_path)
        title = filename.replace('.mp4', '').replace('_', ' ').title()

        # Render video viewer template
        return render_template("student/lesson_video_viewer.html",
                             file_path=file_path,
                             filename=filename,
                             title=title,
                             course_id=course_id,
                             module_index=module_index,
                             lesson_index=lesson_index)

    except Exception as e:
        print(f"Error viewing video: {str(e)}")
        flash("Error loading video.", "danger")
        return redirect(url_for("main.my_courses"))


@main.route('/download_lesson_video/<int:course_id>/<int:module_index>/<int:lesson_index>')
def download_lesson_video(course_id, module_index, lesson_index):
    """Download generated video file"""

    if "user_id" not in session:
        flash("Please log in to download videos.", "warning")
        return redirect(url_for("main.login"))

    try:
        key = f'video_{course_id}_{module_index}_{lesson_index}'
        file_path = session.get(key)

        if file_path and os.path.exists(file_path):
            return send_file(
                file_path,
                as_attachment=True,
                download_name=os.path.basename(file_path),
                mimetype="video/mp4"
            )

        flash("Video not found. Please generate it first.", "danger")
        return redirect(url_for(
            'main.view_my_course',
            course_id=course_id
        ))

    except Exception as e:
        print(f"Error downloading video: {str(e)}")
        traceback.print_exc()

        flash("An error occurred while downloading video.", "danger")
        return redirect(url_for("main.my_courses"))


@main.route('/serve_lesson_video/<int:course_id>/<int:module_index>/<int:lesson_index>')
def serve_lesson_video(course_id, module_index, lesson_index):
    """Serve video file for streaming"""

    if "user_id" not in session:
        return jsonify({"error": "Authentication required"}), 401

    try:
        key = f'video_{course_id}_{module_index}_{lesson_index}'
        video_path = session.get(key)

        if not video_path or not os.path.exists(video_path):
            return jsonify({"error": "Video not found"}), 404

        return send_file(
            video_path,
            mimetype="video/mp4",
            conditional=True  # Enable range requests for streaming
        )

    except Exception as e:
        print(f"Error serving video: {str(e)}")
        return jsonify({"error": str(e)}), 500


# ==================== ASSESSMENT ROUTES ====================

@main.route('/assessment/quiz/<int:course_id>/<int:module_index>/<int:lesson_index>', methods=['GET'])
def get_quiz(course_id, module_index, lesson_index):
    """Get quiz for a specific lesson"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        quiz_data = assessment_service.get_lesson_quiz(course_id, module_index, lesson_index)
        questions = quiz_data.get('questions', []) if isinstance(quiz_data, dict) else quiz_data
        
        return jsonify({
            'status': 'success',
            'data': {
                'course_id': course_id,
                'course_title': quiz_data.get('course_name', f'Course {course_id}') if isinstance(quiz_data, dict) else f'Course {course_id}',
                'module_index': module_index,
                'module_title': quiz_data.get('module_name', f'Module {module_index + 1}') if isinstance(quiz_data, dict) else f'Module {module_index + 1}',
                'lesson_index': lesson_index,
                'lesson_title': quiz_data.get('lesson_name', f'Lesson {lesson_index + 1}') if isinstance(quiz_data, dict) else f'Lesson {lesson_index + 1}',
                'questions': questions,
                'total_questions': len(questions)
            }
        })
    except Exception as e:
        print(f"Error fetching quiz: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to load quiz'}), 500


@main.route('/assessment/quiz/submit', methods=['POST'])
def submit_quiz():
    """Submit quiz responses"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        data = request.get_json()
        user_id = session['user_id']
        course_id = data['course_id']
        module_index = data['module_index']
        lesson_index = data['lesson_index']
        responses = data['responses']
        
        result = assessment_service.submit_quiz_response(
            user_id, course_id, module_index, lesson_index, responses
        )
        
        return jsonify({
            'status': 'success',
            'data': result
        })
    except Exception as e:
        print(f"Error submitting quiz: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to submit quiz'}), 500


@main.route('/assessment/quiz/result/<int:course_id>/<int:module_index>/<int:lesson_index>', methods=['GET'])
def get_quiz_result(course_id, module_index, lesson_index):
    """Get quiz result for a lesson"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        conn = get_db_connection()
        user_id = session['user_id']
        
        result = db_select_one(conn, """
            SELECT total_questions, correct_answers, score_percentage, passed, completed_at
            FROM lesson_quiz_results
            WHERE user_id = %s AND course_id = %s 
            AND module_index = %s AND lesson_index = %s
        """, (user_id, course_id, module_index, lesson_index))
        
        conn.close()
        
        if result:
            return jsonify({
                'status': 'success',
                'data': {
                    'total_questions': result[0],
                    'correct_answers': result[1],
                    'score_percentage': float(result[2]),
                    'passed': bool(result[3]),
                    'completed_at': result[4].isoformat() if result[4] else None
                }
            })
        else:
            return jsonify({'status': 'not_attempted', 'message': 'Quiz not yet attempted'})
    except Exception as e:
        print(f"Error fetching quiz result: {str(e)}")
        return jsonify({'error': 'Failed to load result'}), 500


# ==================== FINAL ASSESSMENT ROUTES ====================

@main.route('/assessment/final/<int:course_id>', methods=['GET'])
def get_final_assessment(course_id):
    """Get final assessment questions"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        assessments = assessment_service.get_final_assessments(course_id)
        
        return jsonify({
            'status': 'success',
            'data': {
                'course_id': course_id,
                'questions': assessments,
                'total_questions': len(assessments)
            }
        })
    except Exception as e:
        print(f"Error fetching final assessment: {str(e)}")
        return jsonify({'error': 'Failed to load assessment'}), 500


@main.route('/assessment/final/submit', methods=['POST'])
def submit_final_assessment():
    """Submit final assessment response"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        data = request.get_json()
        user_id = session['user_id']
        
        result = assessment_service.submit_final_response(
            user_id,
            data['course_id'],
            data['final_assessment_id'],
            data['answer_text']
        )
        
        return jsonify({
            'status': 'success',
            'data': result
        })
    except Exception as e:
        print(f"Error submitting final assessment: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to submit assessment'}), 500


@main.route('/assessment/final/results/<int:course_id>', methods=['GET'])
def get_final_results(course_id):
    """Get final assessment submission and results"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        conn = get_db_connection()
        user_id = session['user_id']
        
        responses = db_select(conn, """
            SELECT sfr.id, sfr.final_assessment_id, sfr.answer_text, sfr.score,
                   sfr.feedback, sfr.is_graded, fa.question_text, fa.max_score
            FROM student_final_responses sfr
            JOIN final_assessments fa ON sfr.final_assessment_id = fa.id
            WHERE sfr.user_id = %s AND sfr.course_id = %s
        """, (user_id, course_id))
        
        conn.close()
        
        results = [
            {
                'id': r[0],
                'assessment_id': r[1],
                'answer': r[2],
                'score': float(r[3]) if r[3] else None,
                'feedback': r[4],
                'is_graded': bool(r[5]),
                'question': r[6],
                'max_score': r[7]
            }
            for r in responses
        ]
        
        return jsonify({
            'status': 'success',
            'data': results
        })
    except Exception as e:
        print(f"Error fetching final results: {str(e)}")
        return jsonify({'error': 'Failed to load results'}), 500


# ==================== LLM-BASED GRADING ENDPOINTS ====================

@main.route('/assessment/llm/grade-saq/<int:response_id>', methods=['POST'])
def llm_grade_saq(response_id):
    """Use LLM to suggest a grade for SAQ response"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        conn = get_db_connection()
        
        # Get response data
        response_data = db_select_one(conn, """
            SELECT ssr.id, ssr.answer_text, saq.question_text, saq.rubric, saq.max_score
            FROM student_saq_responses ssr
            JOIN short_answer_questions saq ON ssr.saq_id = saq.id
            WHERE ssr.id = %s
        """, (response_id,))
        
        conn.close()
        
        if not response_data:
            return jsonify({'error': 'Response not found'}), 404
        
        question_text = response_data[2]
        rubric = response_data[3]
        max_score = response_data[4]
        answer_text = response_data[1]
        
        if rubric:
            rubric = json.loads(rubric) if isinstance(rubric, str) else rubric
        
        # Get LLM grading suggestion
        grading_result = assessment_service.llm_grade_saq_response(
            question_text,
            rubric or {},
            answer_text,
            max_score
        )
        
        return jsonify({
            'status': 'success',
            'data': grading_result,
            'message': 'LLM grading suggestion provided'
        })
    except Exception as e:
        print(f"Error getting LLM grade: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to generate grade'}), 500


@main.route('/assessment/llm/grade-final/<int:response_id>', methods=['POST'])
def llm_grade_final(response_id):
    """Use LLM to suggest a grade for final assessment response"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401
    
    try:
        conn = get_db_connection()
        
        # Get response data
        response_data = db_select_one(conn, """
            SELECT sfr.id, sfr.answer_text, fa.question_text, fa.max_score
            FROM student_final_responses sfr
            JOIN final_assessments fa ON sfr.final_assessment_id = fa.id
            WHERE sfr.id = %s
        """, (response_id,))
        
        conn.close()
        
        if not response_data:
            return jsonify({'error': 'Response not found'}), 404
        
        question_text = response_data[2]
        max_score = response_data[3]
        answer_text = response_data[1]
        
        # Get LLM grading suggestion
        grading_result = assessment_service.llm_grade_final_response(
            question_text,
            answer_text,
            max_score=max_score
        )
        
        return jsonify({
            'status': 'success',
            'data': grading_result,
            'message': 'LLM grading suggestion provided'
        })
    except Exception as e:
        print(f"Error getting LLM grade: {str(e)}")
        traceback.print_exc()
        return jsonify({'error': 'Failed to generate grade'}), 500


@main.route('/assessment/llm/apply-grade/<int:response_id>/<string:response_type>', methods=['POST'])
def apply_llm_grade(response_id, response_type):
    """Apply LLM-suggested grade to response and save it"""
    if 'user_id' not in session:
        return jsonify({'error': 'Unauthorized'}), 401