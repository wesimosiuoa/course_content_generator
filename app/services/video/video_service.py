"""
Video Generation Service
Generates MP4 videos from lesson notes and slides using ElevenLabs TTS and MoviePy
"""

import os
import json
import tempfile
import shutil
from datetime import datetime
from pathlib import Path
from docx import Document
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
import base64
import moviepy.editor as mpy
from dotenv import load_dotenv
load_dotenv()


import os
print("API Key loaded:", bool(os.getenv('ELEVEN_API_KEY')))
# ElevenLabs Configuration
ELEVENLABS_API_KEY = os.getenv('ELEVEN_API_KEY', '')
ELEVENLABS_VOICE_ID = os.getenv('VOICE_ID', 'EXAVITQu4vr4xnSDxMaL') # Default voice

# Base directory for storing generated videos
VIDEO_DIR = os.path.join(os.path.dirname(__file__), 'videos')
MAX_VIDEO_DURATION = 300  # 5 minutes in seconds
TARGET_FPS = 24


def ensure_video_directory():
    """Create video directory if it doesn't exist"""
    if not os.path.exists(VIDEO_DIR):
        os.makedirs(VIDEO_DIR)
    return VIDEO_DIR


def sanitize_filename(text):
    """Remove invalid characters from filename"""
    invalid_chars = '<>:"/\\|?*'
    for char in invalid_chars:
        text = text.replace(char, '_')
    return text.replace(' ', '_').lower()[:50]


def extract_text_from_docx(docx_path):
    """
    Extract all text from a Word document
    
    Args:
        docx_path: Path to DOCX file
        
    Returns:
        str: Extracted text content
    """
    try:
        doc = Document(docx_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""


def extract_slides_from_pptx(pptx_path):
    """
    Extract slides from PPTX with full text content
    
    Args:
        pptx_path: Path to PPTX file
        
    Returns:
        list: List of slide content dictionaries
    """
    try:
        prs = Presentation(pptx_path)
        slides = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_data = {
                'slide_number': slide_num + 1,
                'title': '',
                'text_content': ''
            }
            
            # Extract text from all shapes on the slide
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Check if it's a title (usually the first text shape)
                    if not slide_data['title'] and len(shape.text.strip()) < 100:
                        slide_data['title'] = shape.text.strip()
                    else:
                        text_parts.append(shape.text.strip())
            
            # Combine all text content
            slide_data['text_content'] = ' '.join(text_parts)
            
            slides.append(slide_data)
        
        return slides
    except Exception as e:
        print(f"Error extracting slides from PPTX: {e}")
        return []


def convert_pptx_to_images(pptx_path, output_dir):
    """
    Convert PPTX slides to PNG images with proper text rendering
    
    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save images
        
    Returns:
        list: List of image file paths
    """
    try:
        prs = Presentation(pptx_path)
        image_paths = []
        
        # Try to use a better font, fallback to default if not available
        try:
            font = ImageFont.truetype("arial.ttf", 32)
        except:
            try:
                font = ImageFont.truetype("DejaVuSans.ttf", 32)
            except:
                font = ImageFont.load_default()
        
        for slide_num, slide in enumerate(prs.slides):
            # Create a blank image (1920x1080) with white background
            img = Image.new('RGB', (1920, 1080), color='white')
            draw = ImageDraw.Draw(img)
            
            # Extract and draw text on image with proper line wrapping
            y_position = 80  # Start higher up
            max_width = 1800  # Leave margins
            line_height = 50  # Space between lines
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Split text into words for wrapping
                    words = text.split()
                    lines = []
                    current_line = []
                    
                    for word in words:
                        # Test if adding this word would exceed width
                        test_line = ' '.join(current_line + [word])
                        bbox = draw.textbbox((0, 0), test_line, font=font)
                        text_width = bbox[2] - bbox[0]
                        
                        if text_width <= max_width and current_line:
                            current_line.append(word)
                        else:
                            # Start new line
                            if current_line:
                                lines.append(' '.join(current_line))
                            current_line = [word]
                    
                    # Add remaining line
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    # Draw each line
                    for line in lines:
                        if y_position + line_height > 1000:  # Don't go too low
                            break
                            
                        draw.text((60, y_position), line, fill='black', font=font)
                        y_position += line_height
                    
                    # Add extra space between shapes
                    y_position += 30
            
            # Save image
            img_path = os.path.join(output_dir, f'slide_{slide_num:03d}.png')
            img.save(img_path, 'PNG')
            image_paths.append(img_path)
        
        return image_paths
        
    except Exception as e:
        print(f"Error converting PPTX to images: {e}")
        return []


def generate_tts_audio(text, output_path, voice_id=None):
    """
    Generate audio from text using ElevenLabs TTS API
    """
    try:
        if not ELEVENLABS_API_KEY:
            print("❌ ElevenLabs API key not found in environment")
            return None, 0

        from elevenlabs.client import ElevenLabs
        from elevenlabs import VoiceSettings

        client = ElevenLabs(api_key=ELEVENLABS_API_KEY)
        
        voice_to_use = voice_id or ELEVENLABS_VOICE_ID

        # Generate audio (newer recommended way)
        audio = client.text_to_speech.convert(
            text=text,
            voice_id=voice_to_use,
            model_id="eleven_flash_v2_5",   # or "eleven_turbo_v2_5" for faster
            voice_settings=VoiceSettings(
                stability=0.5,
                similarity_boost=0.75,
                style=0.0,
                use_speaker_boost=True
            )
        )

        # Save audio
        with open(output_path, 'wb') as f:
            for chunk in audio:
                f.write(chunk)

        # Get duration
        audio_clip = mpy.AudioFileClip(output_path)
        duration = audio_clip.duration
        audio_clip.close()

        print(f"✅ TTS generated: {duration:.2f}s")
        return output_path, duration

    except ImportError:
        print("❌ ElevenLabs SDK not installed. Run: pip install elevenlabs")
        return None, 0
    except Exception as e:
        print(f"❌ Error generating TTS audio: {e}")
        return None, 0


def segment_text_for_video(text, num_segments=None):
    """
    Split text into segments for narration
    Each segment should be ~30 seconds of speech
    
    Args:
        text: Full text to segment
        num_segments: Target number of segments (optional)
        
    Returns:
        list: List of text segments
    """
    # Split by sentences
    sentences = text.split('.')
    
    if num_segments:
        # Distribute sentences across segments
        segment_size = max(1, len(sentences) // num_segments)
    else:
        segment_size = 3  # ~3 sentences per segment
    
    segments = []
    current_segment = []
    
    for sentence in sentences:
        sentence = sentence.strip()
        if sentence:
            current_segment.append(sentence)
            
            if len(current_segment) >= segment_size:
                segments.append('. '.join(current_segment) + '.')
                current_segment = []
    
    if current_segment:
        segments.append('. '.join(current_segment) + '.')
    
    return [s.strip() for s in segments if s.strip()]


def generate_video(lesson_notes_path, pptx_path, output_path, course_info=None):
    """
    Main function to generate video from lesson notes and slides
    
    Args:
        lesson_notes_path: Path to DOCX lesson notes
        pptx_path: Path to PPTX presentation
        output_path: Path to save output video
        course_info: Dict with course metadata
        
    Returns:
        dict: Status and metadata of generated video
    """
    temp_dirs = []
    result = {
        'success': False,
        'video_path': None,
        'duration': 0,
        'file_size': 0,
        'slides_count': 0,
        'error': None,
        'message': ''
    }
    
    try:
        ensure_video_directory()
        
        # Create temporary directories
        images_dir = tempfile.mkdtemp()
        audio_dir = tempfile.mkdtemp()
        temp_dirs.extend([images_dir, audio_dir])
        
        # Step 1: Extract content
        print("Extracting lesson notes...")
        lesson_text = extract_text_from_docx(lesson_notes_path)
        
        if not lesson_text:
            raise Exception("No text content found in lesson notes")
        
        print("Extracting slides...")
        slides = extract_slides_from_pptx(pptx_path)
        num_slides = len(slides)
        
        if num_slides == 0:
            raise Exception("No slides found in presentation")
        
        result['slides_count'] = num_slides
        
        # Step 2: Convert slides to images
        print("Converting slides to images...")
        image_paths = convert_pptx_to_images(pptx_path, images_dir)
        
        if not image_paths:
            # Fallback: create placeholder images
            image_paths = []
            for i in range(num_slides):
                img = Image.new('RGB', (1920, 1080), color=(200, 200, 220))
                img_path = os.path.join(images_dir, f'slide_{i:03d}.png')
                img.save(img_path)
                image_paths.append(img_path)
        
        # Step 3: Segment text and generate audio
        print("Generating audio narration...")
        
        # Calculate segments based on slides
        num_segments = min(num_slides, max(1, len(lesson_text.split('.')) // 3))
        text_segments = segment_text_for_video(lesson_text, num_segments)
        
        # Ensure we have at least one segment per slide
        if len(text_segments) < num_slides:
            # Distribute text across more segments
            text_segments = segment_text_for_video(lesson_text, num_slides)
        
        # Generate audio for each segment
        audio_paths = []
        segment_durations = []
        
        for i, segment in enumerate(text_segments[:num_slides]):
            audio_path = os.path.join(audio_dir, f'segment_{i:03d}.mp3')
            seg_path, duration = generate_tts_audio(segment, audio_path)
            
            if seg_path:
                audio_paths.append(seg_path)
                segment_durations.append(duration)
                print(f"  Segment {i+1}: {duration:.2f}s")
            else:
                print(f"  Failed to generate audio for segment {i+1}")
        
        if not audio_paths:
            raise Exception("Failed to generate any audio segments")
        
        # Step 4: Calculate timing
        total_audio_duration = sum(segment_durations)
        print(f"Total audio duration: {total_audio_duration:.2f}s")
        
        if total_audio_duration > MAX_VIDEO_DURATION:
            print(f"Warning: Audio duration ({total_audio_duration:.2f}s) exceeds 5 minutes")
        
        # Calculate slide display duration
        slide_duration = total_audio_duration / len(image_paths) if image_paths else 2
        slide_duration = min(slide_duration, MAX_VIDEO_DURATION / len(image_paths))
        
        # Step 5: Create video clips
        print("Creating video clips...")
        video_clips = []
        
        for i, img_path in enumerate(image_paths):
            # Determine duration for this slide
            if i < len(segment_durations):
                duration = segment_durations[i]
            else:
                duration = slide_duration
            
            # Create clip
            clip = mpy.ImageClip(img_path).set_duration(duration)
            
            # Add audio if available
            if i < len(audio_paths):
                audio_clip = AudioFileClip(audio_paths[i])
                clip = clip.set_audio(audio_clip)
            
            video_clips.append(clip)
        
        # Step 6: Concatenate clips
        print("Assembling video...")
        final_video = mpy.concatenate_videoclips(video_clips)
        
        # Step 7: Write video file
        print(f"Rendering video to {output_path}...")
        final_video.write_videofile(
            output_path,
            fps=TARGET_FPS,
            codec='libx264',
            audio_codec='aac',
            verbose=False,
            logger=None
        )
        
        final_video.close()
        
        # Get file info
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            video_duration = final_video.duration
            
            result['success'] = True
            result['video_path'] = output_path
            result['duration'] = video_duration
            result['file_size'] = file_size
            result['message'] = f"Video generated successfully ({video_duration:.0f}s, {file_size/1024/1024:.2f}MB)"
            
            print(f"✓ Video saved: {output_path}")
        else:
            raise Exception("Video file was not created")
        
    except Exception as e:
        result['error'] = str(e)
        result['message'] = f"Error generating video: {str(e)}"
        print(f"✗ Error: {str(e)}")
    
    finally:
        # Cleanup temporary directories
        for temp_dir in temp_dirs:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
    
    return result


def get_or_create_video(lesson_id, lesson_notes_path, pptx_path, output_filename):
    """
    Get existing video or create a new one
    
    Args:
        lesson_id: ID of the lesson
        lesson_notes_path: Path to lesson notes DOCX
        pptx_path: Path to PPTX presentation
        output_filename: Filename for output video
        
    Returns:
        dict: Video metadata and status
    """
    ensure_video_directory()
    output_path = os.path.join(VIDEO_DIR, output_filename)
    
    # Check if video already exists
    if os.path.exists(output_path):
        file_size = os.path.getsize(output_path)
        return {
            'success': True,
            'video_path': output_path,
            'cached': True,
            'file_size': file_size,
            'message': 'Video loaded from cache'
        }
    
    # Generate new video
    return generate_video(lesson_notes_path, pptx_path, output_path)


def list_generated_videos():
    """
    List all generated videos
    
    Returns:
        list: List of video file dictionaries
    """
    ensure_video_directory()
    videos = []
    
    for filename in os.listdir(VIDEO_DIR):
        if filename.endswith('.mp4'):
            filepath = os.path.join(VIDEO_DIR, filename)
            file_size = os.path.getsize(filepath)
            mtime = os.path.getmtime(filepath)
            
            videos.append({
                'filename': filename,
                'filepath': filepath,
                'file_size': file_size,
                'created_at': datetime.fromtimestamp(mtime).isoformat()
            })
    
    return sorted(videos, key=lambda x: x['created_at'], reverse=True)


def delete_video(video_filename):
    """
    Delete a generated video
    
    Args:
        video_filename: Name of video file to delete
        
    Returns:
        bool: Success status
    """
    ensure_video_directory()
    filepath = os.path.join(VIDEO_DIR, video_filename)
    
    if os.path.exists(filepath) and filepath.startswith(VIDEO_DIR):
        try:
            os.remove(filepath)
            return True
        except Exception as e:
            print(f"Error deleting video: {e}")
            return False
    
    return False
